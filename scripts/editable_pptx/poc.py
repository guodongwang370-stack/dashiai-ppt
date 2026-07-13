"""Offline effect POC for native-text reconstruction."""

from __future__ import annotations

import argparse
import copy
import json
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .masking import changed_outside_mask, composite_masked_edit, make_api_edit_mask
from .renderer import render_editable_pptx
from .scene import EditableScene


@dataclass(frozen=True)
class PocResult:
    pptx_path: Path
    report_path: Path
    clean_plate_path: Path
    internal_mask_path: Path
    api_mask_path: Path


def _load_scene(source: dict[str, Any] | Path | str) -> tuple[dict[str, Any], Path]:
    if isinstance(source, dict):
        return copy.deepcopy(source), Path.cwd()
    path = Path(source).resolve()
    return json.loads(path.read_text(encoding="utf-8")), path.parent


def _resolve_scene_paths(scene: dict[str, Any], base_dir: Path) -> None:
    for key in ("visual_master", "clean_plate", "repair_mask"):
        value = scene.get(key)
        if value and not Path(str(value)).is_absolute():
            scene[key] = str((base_dir / str(value)).resolve())
    for element in scene.get("elements") or []:
        asset = element.get("asset")
        if asset and not Path(str(asset)).is_absolute():
            element["asset"] = str((base_dir / str(asset)).resolve())


def _repair_mask(scene: dict[str, Any]) -> Image.Image:
    canvas = scene.get("canvas") or {}
    width = int(canvas.get("width") or 0)
    height = int(canvas.get("height") or 0)
    repair_mask = scene.get("repair_mask")
    if repair_mask:
        mask = Image.open(str(repair_mask)).convert("L")
        if mask.size != (width, height):
            raise ValueError(f"repair_mask 尺寸必须等于 canvas: {mask.size} != {(width, height)}")
        return mask

    mask = Image.new("L", (width, height), 0)
    draw = ImageDraw.Draw(mask)
    for region in scene.get("repair_regions") or []:
        if not isinstance(region, list) or len(region) != 4:
            raise ValueError("repair_regions 必须包含 [x, y, w, h]")
        x, y, box_width, box_height = (int(value) for value in region)
        if x < 0 or y < 0 or box_width <= 0 or box_height <= 0:
            raise ValueError(f"无效 repair region: {region}")
        if x + box_width > width or y + box_height > height:
            raise ValueError(f"repair region 超出 canvas: {region}")
        draw.rectangle((x, y, x + box_width - 1, y + box_height - 1), fill=255)
    return mask


def run_poc(
    scene_source: dict[str, Any] | Path | str,
    edited_candidate: Path | str,
    output_dir: Path | str,
) -> PocResult:
    scene_data, base_dir = _load_scene(scene_source)
    _resolve_scene_paths(scene_data, base_dir)
    output = Path(output_dir)
    output.mkdir(parents=True, exist_ok=True)

    visual_master_path = Path(str(scene_data.get("visual_master") or ""))
    if not visual_master_path.is_file():
        raise ValueError(f"visual_master 文件不存在: {visual_master_path}")
    visual_master = Image.open(visual_master_path).convert("RGB")
    candidate = Image.open(edited_candidate).convert("RGB")
    mask = _repair_mask(scene_data)
    if visual_master.size != candidate.size or visual_master.size != mask.size:
        raise ValueError(
            f"visual_master、edited candidate 与 mask 尺寸必须一致: "
            f"{visual_master.size}, {candidate.size}, {mask.size}"
        )

    internal_mask_path = output / "internal-mask.png"
    api_mask_path = output / "api-mask.png"
    clean_plate_path = output / "clean-plate.png"
    mask.save(internal_mask_path)
    make_api_edit_mask(mask).save(api_mask_path)
    clean_plate = composite_masked_edit(visual_master, candidate, mask)
    clean_plate.save(clean_plate_path)

    scene_data["clean_plate"] = str(clean_plate_path.resolve())
    editable_scene = EditableScene.from_dict(scene_data)
    pptx_path = render_editable_pptx(editable_scene, output / "editable.pptx")

    shapes = list(Presentation(pptx_path).slides[0].shapes)
    native_text_count = sum(shape.has_text_frame and bool(shape.text) for shape in shapes)
    picture_count = sum(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in shapes)
    expected_text_count = sum(element.type == "native_text" for element in editable_scene.elements)
    outside_changes = changed_outside_mask(visual_master, clean_plate, mask)
    status = "pass" if outside_changes == 0 and native_text_count == expected_text_count else "fail"

    report = {
        "status": status,
        "slide_number": editable_scene.slide_number,
        "outside_mask_changed_pixels": outside_changes,
        "native_text_count": native_text_count,
        "expected_native_text_count": expected_text_count,
        "picture_count": picture_count,
        "visual_master": str(visual_master_path),
        "clean_plate": str(clean_plate_path),
        "pptx": str(pptx_path),
    }
    report_path = output / "quality-report.json"
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")

    return PocResult(
        pptx_path=pptx_path,
        report_path=report_path,
        clean_plate_path=clean_plate_path,
        internal_mask_path=internal_mask_path,
        api_mask_path=api_mask_path,
    )


def main() -> None:
    parser = argparse.ArgumentParser(description="Build an offline editable PPTX POC")
    parser.add_argument("--scene", required=True)
    parser.add_argument("--edited-candidate", required=True)
    parser.add_argument("--output-dir", required=True)
    args = parser.parse_args()
    result = run_poc(args.scene, args.edited_candidate, args.output_dir)
    print(result.report_path)


if __name__ == "__main__":
    main()
