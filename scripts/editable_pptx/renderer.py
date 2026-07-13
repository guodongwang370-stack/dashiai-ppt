"""Render an EditableScene as native PowerPoint objects."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from .scene import EditableScene, SceneElement


ALIGNMENTS = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}
VERTICAL_ALIGNMENTS = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}
SHAPE_TYPES = {
    "rectangle": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    "rounded_rectangle": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
    "star_5": MSO_AUTO_SHAPE_TYPE.STAR_5_POINT,
}


def _hex_color(value: str) -> RGBColor:
    normalized = value.strip().lstrip("#")
    if len(normalized) != 6:
        raise ValueError(f"无效颜色: {value}")
    return RGBColor.from_string(normalized.upper())


def _slide_bbox(prs: Presentation, scene: EditableScene, element: SceneElement) -> tuple[int, int, int, int]:
    x, y, width, height = element.bbox_px
    return (
        int(x / scene.canvas_width * prs.slide_width),
        int(y / scene.canvas_height * prs.slide_height),
        int(width / scene.canvas_width * prs.slide_width),
        int(height / scene.canvas_height * prs.slide_height),
    )


def _add_text(prs: Presentation, slide, scene: EditableScene, element: SceneElement) -> None:
    left, top, width, height = _slide_bbox(prs, scene, element)
    shape = slide.shapes.add_textbox(left, top, width, height)
    shape.name = element.id
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    frame.vertical_anchor = VERTICAL_ALIGNMENTS.get(
        str(element.style.get("vertical_align", "middle")), MSO_ANCHOR.MIDDLE
    )

    paragraph = frame.paragraphs[0]
    paragraph.alignment = ALIGNMENTS.get(str(element.style.get("align", "left")), PP_ALIGN.LEFT)
    run = paragraph.add_run()
    run.text = element.content
    font = run.font
    font.name = str(element.style.get("font_face", "Noto Sans CJK SC"))
    font.size = Pt(float(element.style.get("font_size_pt", 24)))
    font.bold = int(element.style.get("font_weight", 400)) >= 600
    font.color.rgb = _hex_color(str(element.style.get("color", "#000000")))


def _add_shape(prs: Presentation, slide, scene: EditableScene, element: SceneElement) -> None:
    shape_name = str(element.style.get("shape"))
    if shape_name == "line":
        _add_connector(prs, slide, scene, element)
        return
    left, top, width, height = _slide_bbox(prs, scene, element)
    shape = slide.shapes.add_shape(SHAPE_TYPES[shape_name], left, top, width, height)
    shape.name = element.id
    fill_color = element.style.get("fill")
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_color(str(fill_color))
        shape.fill.transparency = float(element.style.get("fill_transparency", 0))
    else:
        shape.fill.background()
    line_color = element.style.get("line")
    if line_color:
        shape.line.color.rgb = _hex_color(str(line_color))
        shape.line.transparency = float(element.style.get("line_transparency", 0))
        shape.line.width = Pt(float(element.style.get("line_width_pt", 1)))
    else:
        shape.line.fill.background()
    shape.rotation = float(element.style.get("rotation", 0))


def _add_connector(prs: Presentation, slide, scene: EditableScene, element: SceneElement) -> None:
    left, top, width, height = _slide_bbox(prs, scene, element)
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        left,
        top,
        left + width,
        top + height,
    )
    connector.name = element.id
    connector.line.color.rgb = _hex_color(str(element.style.get("line", "#000000")))
    connector.line.width = Pt(float(element.style.get("line_width_pt", 1)))
    connector.line.transparency = float(element.style.get("line_transparency", 0))
    if element.style.get("end_arrow"):
        line_properties = connector.line._get_or_add_ln()
        tail_end = OxmlElement("a:tailEnd")
        tail_end.set("type", "triangle")
        line_properties.append(tail_end)


def _render_element(prs: Presentation, slide, scene: EditableScene, element: SceneElement) -> None:
    if element.type == "image_layer":
        left, top, width, height = _slide_bbox(prs, scene, element)
        picture = slide.shapes.add_picture(str(element.asset), left, top, width=width, height=height)
        picture.name = element.id
    elif element.type == "native_text":
        _add_text(prs, slide, scene, element)
    elif element.type == "native_shape":
        _add_shape(prs, slide, scene, element)
    elif element.type == "connector":
        _add_connector(prs, slide, scene, element)


def add_editable_slide(prs: Presentation, scene: EditableScene):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    clean_plate = slide.shapes.add_picture(
        str(scene.clean_plate),
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )
    clean_plate.name = "clean_plate"
    for element in sorted(scene.elements, key=lambda item: item.z_index):
        _render_element(prs, slide, scene, element)
    return slide


def render_editable_deck(scenes, output_path: Path | str) -> Path:
    ordered = sorted(scenes, key=lambda scene: scene.slide_number)
    if not ordered:
        raise ValueError("至少需要一个 editable scene")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for scene in ordered:
        add_editable_slide(prs, scene)
    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    return path


def render_editable_pptx(scene: EditableScene, output_path: Path | str) -> Path:
    return render_editable_deck([scene], output_path)
