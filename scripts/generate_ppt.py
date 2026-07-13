#!/usr/bin/env python3
"""
PPT Generator - Generate PPT slide images using OpenAI gpt-image-2 (Images API).

Generates 16:9 slide images from a slide plan + style template, saves
structured slide_spec metadata for precise editing, and packages a .pptx.
"""

import argparse
import copy
import json
import os
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

# =============================================================================
# Constants
# =============================================================================

OUTPUT_BASE_DIR = "outputs"
SLIDE_TRACE_WIDTH = 1536
SLIDE_TRACE_HEIGHT = 864
EXTERNAL_IMAGE_TYPES = {"external_image", "external_image_placeholder", "image_overlay"}
REFERENCE_ONLY_IMAGE_TYPES = {"image_reference", "visual_reference", "reference_image", "style_reference_image"}
REFERENCE_ONLY_MODES = {"reference", "image_reference", "image-to-image", "image_to_image", "redraw", "blend", "inspiration"}

SCRIPT_DIR = Path(__file__).parent
SKILL_ROOT = SCRIPT_DIR.parent
CWD = Path.cwd()


# =============================================================================
# Environment Configuration
# =============================================================================

SKILL_PREFIX = "gpt-image2-ppt"

SKILL_ENV_MAP = {
    "OPENAI_API_KEY":        "OPENAI_API_KEY",
    "OPENAI_BASE_URL":       "OPENAI_BASE_URL",
    "GPT_IMAGE_MODEL_NAME":  "GPT_IMAGE_MODEL_NAME",
    "GPT_IMAGE_QUALITY":     "GPT_IMAGE_QUALITY",
    "GPT_IMAGE_BACKEND":     "GPT_IMAGE_BACKEND",
    "GPT_IMAGE_CONCURRENCY": "GPT_IMAGE_CONCURRENCY",
    "VISION_BASE_URL":       "VISION_BASE_URL",
    "VISION_API_KEY":        "VISION_API_KEY",
    "VISION_MODEL_NAME":     "VISION_MODEL_NAME",
}


def _load_platform_env() -> bool:
    """Map platform-injected 'skill-prefix_VAR' env vars to unprefixed names.

    The platform (MedAgent / Codex / OpenClaw) sets env vars like
        gpt-image2-ppt_OPENAI_API_KEY=sk-...
    but the script expects plain OPENAI_API_KEY.  This function detects
    the prefixed form and maps them transparently into os.environ.
    Returns True if any prefixed var was found.
    """
    found = False
    for base_name in SKILL_ENV_MAP:
        prefixed = f"{SKILL_PREFIX}_{base_name}"
        value = os.environ.get(prefixed)
        if value and not os.environ.get(base_name):
            os.environ[base_name] = value
            found = True
    return found


def _load_scoped_env_files() -> List[str]:
    """Load env files from explicit or skill-owned locations only.

    This intentionally does not read CWD/.env or walk parent directories. Skill
    credentials should come from the agent framework / system environment first;
    .env is only a standalone CLI fallback.
    """
    candidates: List[Path] = []
    explicit = os.environ.get("GPT_IMAGE2_PPT_ENV")
    if explicit:
        candidates.append(Path(explicit).expanduser())
    candidates.extend([
        SKILL_ROOT / ".env",
        Path.home() / ".codex/skills/gpt-image2-ppt-skills/.env",
        Path.home() / ".claude/skills/gpt-image2-ppt-skills/.env",
        Path.home() / "skills/gpt-image2-ppt/.env",
        Path.home() / "skills/gpt-image2-ppt-skills/.env",
    ])

    try:
        from dotenv import load_dotenv
    except ImportError:
        return []

    loaded: List[str] = []
    seen = set()
    for path in candidates:
        resolved = str(path.resolve()) if path.exists() else str(path)
        if resolved in seen or not path.is_file():
            continue
        seen.add(resolved)
        if load_dotenv(path, override=False):
            loaded.append(resolved)
    return loaded


def load_skill_env() -> None:
    """Load environment with safe precedence.

    Existing process environment wins. Platform-prefixed variables are mapped
    next, then explicit / skill-owned .env files fill in missing values.
    """
    _load_platform_env()
    _load_scoped_env_files()
    # Provider-specific convenience aliases. If the user explicitly provides
    # JULING_GPT_IMAGE2_* for this skill, prefer that image channel.
    juling_base = os.environ.get("JULING_GPT_IMAGE2_BASE_URL")
    juling_key = os.environ.get("JULING_GPT_IMAGE2_API_KEY")
    if juling_base:
        os.environ["OPENAI_BASE_URL"] = juling_base
    if juling_key:
        os.environ["OPENAI_API_KEY"] = juling_key

# =============================================================================
# Session & Metadata Management
# =============================================================================

METADATA_FILENAME = "metadata.json"


def _find_sessions(base_dir: Optional[str] = None) -> List[Dict[str, Any]]:
    """Discover all generation sessions (directories containing metadata.json).

    Returns a list of dicts with 'timestamp', 'dir', 'title', 'slides'.
    """
    if base_dir is None:
        base_dir = str(CWD / OUTPUT_BASE_DIR)
    sessions: List[Dict[str, Any]] = []
    if not os.path.isdir(base_dir):
        return sessions
    for entry in sorted(os.listdir(base_dir), reverse=True):
        session_dir = os.path.join(base_dir, entry)
        meta_path = os.path.join(session_dir, METADATA_FILENAME)
        if os.path.isfile(meta_path):
            try:
                with open(meta_path, "r", encoding="utf-8") as f:
                    meta = json.load(f)
                sessions.append({
                    "timestamp": entry,
                    "dir": session_dir,
                    "title": meta.get("title", "Untitled"),
                    "slide_count": len(meta.get("slides", {})),
                })
            except Exception:
                sessions.append({
                    "timestamp": entry,
                    "dir": session_dir,
                    "title": "(corrupt metadata)",
                    "slide_count": 0,
                })
    return sessions


def _resolve_session(session_id: str) -> str:
    """Resolve a session timestamp or path to the session directory.

    Accepts: full path, relative path, or timestamp in OUTPUT_BASE_DIR.
    """
    # Absolute or relative path
    if os.path.isdir(session_id) and os.path.isfile(os.path.join(session_id, METADATA_FILENAME)):
        return os.path.abspath(session_id)
    # Try under OUTPUT_BASE_DIR
    candidate = str(CWD / OUTPUT_BASE_DIR / session_id)
    if os.path.isdir(candidate) and os.path.isfile(os.path.join(candidate, METADATA_FILENAME)):
        return candidate
    # Try with partial match (prefix)
    sessions = _find_sessions()
    for s in sessions:
        if s["timestamp"].startswith(session_id):
            return s["dir"]
    print(f"[X] Session not found: {session_id}")
    print(f"    Tried: {session_id}")
    if candidate != session_id:
        print(f"    Tried: {candidate}")
    sys.exit(1)


def _load_metadata(session_dir: str) -> Dict[str, Any]:
    """Load metadata.json from a session directory."""
    meta_path = os.path.join(session_dir, METADATA_FILENAME)
    if not os.path.isfile(meta_path):
        print(f"[X] No metadata.json found in {session_dir}")
        sys.exit(1)
    with open(meta_path, "r", encoding="utf-8") as f:
        return json.load(f)


def _save_metadata(metadata: Dict[str, Any], session_dir: str) -> None:
    """Atomically save metadata.json to a session directory.

    Writes to a temp file first, then renames to avoid corruption
    on crashes or disk-full conditions.
    """
    import tempfile as _tempfile
    meta_path = os.path.join(session_dir, METADATA_FILENAME)
    fd, tmp_path = _tempfile.mkstemp(
        suffix=".json", prefix=".metadata_", dir=session_dir
    )
    try:
        with os.fdopen(fd, "w", encoding="utf-8") as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)
        os.replace(tmp_path, meta_path)  # atomic on POSIX + modern Windows
    except Exception:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass
        raise


def _collect_slide_numbers(metadata: Dict[str, Any]) -> List[int]:
    """Extract sorted slide numbers from metadata slides dict."""
    slides_dict = metadata.get("slides", {})
    ordered: List[int] = []
    seen = set()
    for raw in metadata.get("slide_order", []) or []:
        try:
            n = int(raw)
        except (TypeError, ValueError):
            continue
        if str(n) in slides_dict and n not in seen:
            ordered.append(n)
            seen.add(n)

    nums: List[int] = []
    for k in slides_dict.keys():
        try:
            n = int(k)
        except ValueError:
            continue
        if n not in seen:
            nums.append(n)
    nums.sort()
    return ordered + nums


def _get_latest_slide_spec(metadata: Dict[str, Any], slide_number: int) -> Dict[str, Any]:
    """Get the current slide_spec for a given slide number."""
    slide_key = str(slide_number)
    slide_data = metadata.get("slides", {}).get(slide_key)
    if not slide_data:
        return {}
    current_version = slide_data.get("current_version", 1)
    for v in slide_data.get("versions", []):
        if v.get("version") == current_version:
            return v.get("spec", {})
    return {}


def _get_version_info(metadata: Dict[str, Any], slide_number: int, version: int) -> Optional[Dict[str, Any]]:
    """Get a specific version entry for a slide."""
    slide_key = str(slide_number)
    slide_data = metadata.get("slides", {}).get(slide_key)
    if not slide_data:
        return None
    for v in slide_data.get("versions", []):
        if v.get("version") == version:
            return v
    return None


def _init_slide_metadata(
    slide_number: int,
    page_type: str,
    initial_spec: Dict[str, Any],
    prompt_file: str,
    image_path: str,
) -> Dict[str, Any]:
    """Create the initial metadata entry for a slide."""
    return {
        "slide_number": slide_number,
        "page_type": page_type,
        "current_version": 1,
        "image_snapshot": image_path,
        "versions": [
            {
                "version": 1,
                "action": "generate",
                "spec": initial_spec,
                "prompt_file": prompt_file,
                "image_snapshot": image_path,
            }
        ],
    }


def _add_slide_version(
    slide_data: Dict[str, Any],
    new_version: int,
    spec: Dict[str, Any],
    action: str,
    image_snapshot: str,
    prompt_file: str,
    edit_instruction: str = "",
    reference_version: Optional[int] = None,
) -> None:
    """Append a new version entry to slide metadata."""
    entry: Dict[str, Any] = {
        "version": new_version,
        "action": action,
        "spec": spec,
        "prompt_file": prompt_file,
        "image_snapshot": image_snapshot,
    }
    if edit_instruction:
        entry["edit_instruction"] = edit_instruction
    if reference_version is not None:
        entry["reference_version"] = reference_version
    slide_data.setdefault("versions", []).append(entry)
    slide_data["current_version"] = new_version
    slide_data["image_snapshot"] = image_snapshot


def _stabilize_version_snapshots(slide_data: Dict[str, Any], slide_number: int, images_dir: str) -> None:
    """Point version history entries at immutable versioned image files when present."""
    for version_entry in slide_data.get("versions", []):
        version = version_entry.get("version")
        if not isinstance(version, int):
            continue
        versioned_rel = f"images/slide-{slide_number:02d}_v{version:04d}.png"
        versioned_abs = os.path.join(images_dir, f"slide-{slide_number:02d}_v{version:04d}.png")
        if os.path.isfile(versioned_abs):
            version_entry["image_snapshot"] = versioned_rel


# =============================================================================
# Slide Spec Helpers
# =============================================================================

def apply_spec_updates(spec: Dict[str, Any], element_updates: Dict[str, Any]) -> Dict[str, Any]:
    """Apply element-level updates to a slide_spec, returning the updated spec.

    element_updates maps element_id to a dict of key-value changes.
    e.g. {"subtitle": {"content": "新内容"}, "title": {"color": "#ff0000"}}
    """
    import copy
    updated = copy.deepcopy(spec)
    elements = updated.setdefault("elements", {})
    derived_external_keys = {
        "computed_bbox",
        "placement_region",
        "asset_size",
        "asset_ratio",
        "tailored_to_asset",
        "auto_layout_reason",
    }
    geometry_keys = {
        "position",
        "bbox",
        "source",
        "path",
        "asset",
        "fit",
        "slot",
        "slot_strategy",
        "tailor_to_asset",
        "anchor",
        "padding",
        "bleed",
        "layout_intent",
        "auto_layout",
    }
    for elem_id, changes in element_updates.items():
        if elem_id in elements:
            elem = elements[elem_id]
            if (
                isinstance(elem, dict)
                and str(elem.get("type", "")).lower() in {"external_image", "external_image_placeholder", "image_overlay"}
                and any(k in changes for k in geometry_keys)
            ):
                for k in derived_external_keys:
                    elem.pop(k, None)
            elements[elem_id].update(changes)
        else:
            elements[elem_id] = changes
    return updated


def construct_edit_prompt(spec: Dict[str, Any], element_updates: Dict[str, Any]) -> str:
    """Construct a natural-language edit prompt from the old spec and element changes.

    The prompt instructs gpt-image-2 to change only specific elements
    while keeping everything else identical to the reference image.
    """
    parts = ["在参考图基础上，只修改以下内容，保持其他所有元素（背景、装饰、布局、字体、颜色、大小）完全不变："]
    elements = spec.get("elements", {})

    for elem_id, changes in element_updates.items():
        old_elem = elements.get(elem_id, {})
        position = old_elem.get("position", "相应位置")
        etype = old_elem.get("type", "元素")
        old_content = old_elem.get("content", "")
        new_content = changes.get("content", old_content)

        if "content" in changes and old_content and new_content != old_content:
            parts.append(
                f"将{position}{etype}的文字从「{old_content}」改为「{new_content}」"
            )
        elif "color" in changes:
            old_color = old_elem.get("color", "")
            new_color = changes.get("color", "")
            parts.append(
                f"将{position}{etype}的颜色从{old_color}改为{new_color}"
            )
        elif "style" in changes:
            parts.append(
                f"将{position}{etype}的样式改为{changes['style']}"
            )
        else:
            for k, v in changes.items():
                parts.append(f"修改{position}{etype}的{k}为{v}")

    return "\n".join(parts)


def construct_external_slot_repair_prompt(spec: Dict[str, Any], element_updates: Optional[Dict[str, Any]] = None) -> str:
    """Prompt for repairing a slide whose content collides with external-image slots."""
    slot_ids = []
    for elem_id, elem in (spec.get("elements") or {}).items():
        if isinstance(elem, dict) and str(elem.get("type", "")).lower() in {
            "external_image",
            "external_image_placeholder",
            "image_overlay",
        }:
            slot_ids.append(str(elem_id))
    slot_text = "、".join(slot_ids) if slot_ids else "外部真实图片槽位"
    parts = [
        "请基于当前幻灯片图进行版式修复，保持整体风格、配色、字体气质、背景质感和主要视觉层级不变。",
        f"另一张参考图是空间定位图，细角标标出的区域是后续会贴入真实图片的覆盖区（{slot_text}），不是页面设计元素。",
        "必须把所有标题、正文、数字、图标、线条、装饰图形、照片主体和重要视觉元素移出这些角标区域。",
        "允许为了避让槽位而小幅移动、缩短、换行、重排或重新对齐文字和图形；优先保持页面清爽和可读。",
        "修复后仍要保持模板原本的结构节奏：图文分栏、上下区块、主视觉方向和留白比例要自然，不能为了避让产生大片无意义空白、断裂分割线、挤在角落的文字或失衡构图。",
        "角标区域下方只生成自然连续的背景、轻微纹理或非常轻的抽象装饰。",
        "严禁在角标区域生成白色/浅色矩形、空卡片、图片占位框、边框、阴影、透明洞或任何可见预留块。",
        "不要复制、描摹、强化、美化或保留空间定位图里的角标；最终页面里应看不到任何占位痕迹。",
        "不要生成或仿造真实图片内容，真实图片会在 PPTX 打包阶段由代码按同一坐标贴入。",
    ]
    if element_updates:
        parts.append("同时应用以下 slide_spec 位置/内容更新，并以更新后的空间定位图为准：")
        parts.append(json.dumps(element_updates, ensure_ascii=False, indent=2))
    return "\n".join(parts)


def _format_external_slots_constraint(external_slots: List[Dict[str, Any]]) -> str:
    """Prompt block for pages where real external images are overlaid later."""
    if not external_slots:
        return ""
    slot_lines = []
    for slot in external_slots:
        tailored = "已按真实素材宽高比量体裁衣" if slot.get("tailored_to_asset") else "固定槽位"
        asset_info = ""
        if slot.get("asset_size") and slot.get("asset_ratio"):
            asset_info = f"，素材尺寸={slot.get('asset_size')}，素材比例={slot.get('asset_ratio')}"
        if slot.get("asset_class"):
            asset_info += f"，素材类别={slot.get('asset_class')}，细节级别={slot.get('detail_level')}"
        slot_lines.append(
            f"- {slot['id']}（{tailored}{asset_info}）：对应参考图中的一组细角标定位区。"
        )
        if slot.get("auto_layout_reason"):
            slot_lines.append(f"  版式规划: {slot.get('auto_layout_reason')}")
    return (
        "\n\n【外部真实图片槽位 - 强约束】\n"
        "在绘制页面前，先根据本页文字量、真实图片数量、真实图片比例和参考图角标位置做整体排版规划。"
        "不要照搬风格模板中原有的照片区、图片框或主体摄影构图；它们只提供风格气质，不提供本页最终图片内容。\n"
        "参考图中的细角标只标出后续会由代码贴入真实图片的覆盖区，"
        "它们不是页面设计元素，也不是图片框、卡片、遮罩或留白块。\n"
        + "\n".join(slot_lines)
        + "\n请只避让这些定位区里的关键信息：不要把标题、正文、数字、图标、人物或主体照片放进这些区域。"
        + "定位区底下的背景必须自然连续，可以有同页背景色、纸张纹理或轻微装饰延续过去。"
        + "严禁在定位区生成白色/浅色矩形、空卡片、图片占位框、边框、阴影、透明洞或任何可见预留块。"
        + "如果页面布局或风格模板提到主体照片、产品照、人物照、场景摄影、photo crop、image crop、图片区或照片区，"
        + "这些角色都将由代码后贴真实图片完成，模型绝对不要自行生成、仿造或重绘这些照片内容。"
        + "本页可以保留本风格的文字、细线、背景纹理、柔和光影和少量抽象装饰，但不要生成额外摄影主体。"
        + "如果风格模板的照片/图像构图与外部真实图片槽位冲突，以外部槽位、文字可读性和页面清爽为最高优先级。"
        + "不要复制、描摹、强化或美化参考图里的角标；最终页面里应看不到任何占位痕迹。"
    )


def _compact_text_for_asset_first_layout(content: Any, max_items: int = 3, max_chars_per_item: int = 24) -> str:
    text = str(content or "").strip()
    if not text:
        return ""
    parts = [
        p.strip()
        for p in re.split(r"[\n。；;.!！？?]+", text)
        if p.strip()
    ]
    compact: List[str] = []
    for part in parts:
        part = re.sub(r"\s+", " ", part).strip()
        if len(part) > max_chars_per_item:
            part = part[: max_chars_per_item - 1].rstrip() + "…"
        compact.append(part)
        if len(compact) >= max_items:
            break
    return "\n".join(f"要点：{p}" for p in compact)


def _wrap_title_for_asset_first_layout(content: Any, max_chars_per_line: int = 12, max_lines: int = 3) -> str:
    text = re.sub(r"\s+", " ", str(content or "").strip())
    if not text:
        return ""
    lines: List[str] = []
    current = ""
    for ch in text:
        current += ch
        if len(current) >= max_chars_per_line or ch in "：:，,、/":
            lines.append(current.strip())
            current = ""
            if len(lines) >= max_lines:
                break
    if current and len(lines) < max_lines:
        lines.append(current.strip())
    return "\n".join(line for line in lines if line)


def _format_asset_first_layout_constraint(high_detail_slots: List[Dict[str, Any]]) -> str:
    if not high_detail_slots:
        return ""
    slot_names = "、".join(str(slot.get("id", "真实图")) for slot in high_detail_slots)
    return (
        "\n\n【高细节真实图主视觉页 - 强约束】\n"
        f"本页包含需要阅读细节的真实素材（{slot_names}），例如论文截图、表格、架构图、代码、公式或曲线图。"
        "这些真实图不是装饰缩略图，而是页面主视觉，必须预留大面积可读区域。\n"
        "版式策略：采用“左侧窄文字栏 + 右侧/下方大素材面板”结构；文字区必须明显避开角标素材区。"
        "当素材是标准比例或竖图时，标题和摘要都收进左侧窄栏，标题可以换行变短，但不要横跨到右侧素材面板。\n"
        "当素材是超宽图时，标题和摘要收进上方/左上文字区，不要下探到下方素材面板。\n"
        "文字策略：正文最多 2-3 条短要点，每条尽量不超过 18 个中文字符；不要把原文长段落完整铺进页面。"
        "如果文字和真实图可读性冲突，优先牺牲文字数量，保留大图可读性。\n"
        "视觉策略：素材区底下只生成连续背景，不生成白卡、图片框、假截图、假论文页或装饰块。"
    )


def _strip_generated_image_directives(text: str) -> str:
    """Remove photo/image-generation directives from a prompt for external-image pages."""
    photo_terms = (
        "image",
        "photo",
        "photograph",
        "photography",
        "crop",
        "full-bleed",
        "imagery",
        "照片",
        "摄影",
        "图片",
        "图像",
        "图片区",
        "图像区",
        "照片区",
    )
    keep_markers = (
        "外部真实图片",
        "真实图片会",
        "后续会由代码贴入",
        "reference skeleton",
        "asset-skeleton",
    )
    lines = []
    for line in text.splitlines():
        stripped = line.strip()
        lowered = stripped.lower()
        if any(marker in stripped for marker in keep_markers):
            lines.append(line)
            continue
        if any(term in lowered for term in photo_terms):
            continue
        lines.append(line)
    return "\n".join(lines).strip()


def adapt_template_prompt_for_external_slots(prompt: str, external_slots: List[Dict[str, Any]]) -> str:
    """Make a template-derived prompt safe for pages with real image overlays."""
    if not external_slots:
        return prompt
    prompt = _strip_generated_image_directives(prompt)
    return (
        prompt
        + "\n\n【模板冲突消解 - 最高优先级】\n"
        + "本页不是按模板照片区直接重画照片；模板只提供配色、字体、网格、节奏和装饰语言。"
        + "真实图片的位置、数量和比例由外部图片槽位与参考图角标决定。"
        + "请先根据本页文字量和真实图片槽位重新组织版式：文字必须绕开槽位，图形和装饰不得侵入槽位。"
        + "如果模板布局描述中仍隐含照片、图片区、全幅图、裁切图或视觉主体图，请把这些指令视为已被外部真实图片替代。"
        + "最终生成图中不要出现额外照片、假图、图片框、占位卡片或为了图片区而留下的白色/浅色底块。"
    )


def generate_prompt_from_spec(
    style_template: str,
    slide_spec: Dict[str, Any],
    page_type: str,
    slide_number: int,
    total_slides: int,
    output_dir: Optional[str] = None,
) -> str:
    """Generate a detailed prompt from a structured slide_spec.

    Uses element-level descriptions (type, content, position, style, color)
    to construct a precise prompt that gpt-image-2 can follow.
    """
    elements = slide_spec.get("elements", {})
    layout = slide_spec.get("layout", "")
    external_slots = _collect_external_image_slots(slide_spec)
    if output_dir:
        external_slots = _slots_with_real_sources(external_slots, output_dir)
    high_detail_slot_ids = {
        str(slot.get("id"))
        for slot in external_slots
        if _external_asset_detail_level(slide_spec, str(slot.get("id")), (elements.get(str(slot.get("id"))) or {})) in {"high", "critical"}
    }
    effective_style_template = (
        _adapt_style_template_for_external_slots(style_template)
        if external_slots
        else style_template
    )

    # Build element descriptions
    element_lines = []
    for elem_id, elem in elements.items():
        etype = elem.get("type", "unknown")
        if str(etype).lower() in EXTERNAL_IMAGE_TYPES or str(etype).lower() in REFERENCE_ONLY_IMAGE_TYPES:
            continue
        content = _element_text_content(elem)
        etype_lower = str(etype).lower()
        if high_detail_slot_ids and etype_lower in {"title", "headline", "heading"}:
            content = _wrap_title_for_asset_first_layout(content)
        elif high_detail_slot_ids and etype_lower in {"body", "paragraph", "text", "bullet", "bullets"}:
            content = _compact_text_for_asset_first_layout(content, max_items=2, max_chars_per_item=16)
        position = elem.get("position", "")
        if high_detail_slot_ids and not position:
            if etype_lower in {"title", "headline", "heading"}:
                position = "左侧窄文字栏上方，标题必须换成短行，不要横向延伸到右侧素材面板"
            elif etype_lower in {"body", "paragraph", "text", "bullet", "bullets"}:
                position = "左侧窄文字栏内、标题下方，只保留短要点，不要进入右侧素材面板"
        style_hint = elem.get("style", "")
        color = elem.get("color", "")
        description = elem.get("description", "")

        if description:
            element_lines.append(f"- {elem_id}（{etype}）: {description}")
            continue

        desc_parts = [f"- {etype}「{content}」"]
        if position:
            desc_parts.append(f"位置: {position}")
        if style_hint:
            desc_parts.append(f"样式: {style_hint}")
        if color:
            desc_parts.append(f"颜色: {color}")
        element_lines.append("，".join(desc_parts))

    elements_text = "\n".join(element_lines)
    if external_slots:
        elements_text += _format_external_slots_constraint(external_slots)
        high_detail_slots = [
            slot for slot in external_slots
            if str(slot.get("id")) in high_detail_slot_ids
        ]
        elements_text += _format_asset_first_layout_constraint(high_detail_slots)

    # Page type hint
    is_cover = page_type == "cover" or slide_number == 1
    is_data = page_type == "data" or slide_number == total_slides
    if is_cover:
        label = "封面页（cover）"
        hint = "标题/副标题处理为视觉焦点，按本风格的封面构图规范处理。"
    elif is_data:
        label = "数据页（data）"
        hint = "突出关键数字、对比或结论；按本风格的数据/总结构图规范处理。"
    else:
        label = "内容页（content）"
        hint = "把要点按本风格的内容构图规范结构化呈现，注意层级、对齐、留白。"

    if external_slots:
        if high_detail_slot_ids:
            layout_line = (
                "\n页面布局: 高细节真实图主视觉页。先围绕参考图角标区域建立大素材面板，"
                "再把标题和 2-3 条短摘要放入剩余窄文字栏；模型只生成背景、短文字、细线、编号和抽象装饰，"
                "不生成任何照片、截图或图片区。"
            )
        else:
            layout_line = (
                "\n页面布局: 按文字元素的位置组织版式；参考图角标区域只作为后续真实图片覆盖区，"
                "模型只生成背景、文字、细线、编号和抽象装饰，不生成任何照片或图片区。"
            )
    else:
        layout_line = f"\n页面布局: {layout}" if layout else ""

    return (
        effective_style_template
        + "\n\n---\n\n"
        + f"现在请生成本组中的【{label}】，{hint}{layout_line}"
        + "\n\n本页各元素的精确描述（请严格按以下布局、位置、样式生成）：\n\n"
        + elements_text
        + LANGUAGE_FONT_RULE
    )


def _stringify_content_value(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, (list, tuple)):
        return "；".join(str(x).strip() for x in value if str(x).strip())
    if isinstance(value, dict):
        return json.dumps(value, ensure_ascii=False)
    return str(value).strip()


def _element_text_content(elem: Dict[str, Any], *, include_description: bool = False) -> str:
    """Return human-facing text from a slide_spec element.

    SKILL.md documents both `content` and `heading`/`body` element shapes.
    Prompt generation and text-density heuristics must serialize either shape,
    otherwise card-like specs become empty text guides.
    """
    content = _stringify_content_value(elem.get("content"))
    if content:
        return content

    parts: List[str] = []
    for key in ("heading", "title", "label"):
        value = _stringify_content_value(elem.get(key))
        if value:
            parts.append(value)
            break
    for key in ("body", "text"):
        value = _stringify_content_value(elem.get(key))
        if value:
            parts.append(value)
            break

    items = elem.get("items")
    if items is None:
        items = elem.get("bullets")
    item_text = _stringify_content_value(items)
    if item_text:
        parts.append(item_text)

    if parts:
        return "：".join(parts[:2]) + ("；" + "；".join(parts[2:]) if len(parts) > 2 else "")

    if include_description:
        return _stringify_content_value(elem.get("description"))
    return ""


def _adapt_style_template_for_external_slots(style_template: str) -> str:
    """Remove generated-photo directives when real images are overlaid later."""
    skipped_sections = {
        "封面页构图",
        "内容页构图",
        "布局系统",
        "图片处理",
        "章节页构图",
        "收尾页构图",
    }
    photo_terms = (
        "image",
        "photo",
        "photograph",
        "photography",
        "crop",
        "full-bleed",
        "imagery",
        "照片",
        "摄影",
        "图片",
        "图像",
    )
    lines = []
    skip = False
    for line in style_template.splitlines():
        stripped = line.strip()
        section_match = re.match(r"^【(.+?)】$", stripped)
        if section_match:
            skip = section_match.group(1) in skipped_sections
            if skip:
                continue
        if skip:
            continue
        if any(term in stripped.lower() for term in photo_terms):
            continue
        lines.append(line)

    override = """

【外部真实图片后贴页覆盖规则 - 最高优先级】
- 本页真实图片会在生成后由代码精确贴入；模型不要生成任何照片、图片裁切、产品照、人物照、场景摄影或类似主体。
- 保留所选风格本身的配色、字体气质、结构线、留白、背景质感和抽象装饰语言。
- 参考图角标区域下方只允许连续背景、轻微纹理或非常轻的抽象装饰；不要出现白色占位块、图片框、卡片、边框、阴影或透明洞。
- 如需要视觉重心，用字号、细线、编号、留白、色块或轻微纹理实现，不要用生成照片实现。
"""
    return "\n".join(lines).strip() + override


# =============================================================================
# Style Template
# =============================================================================

def load_style_template(style_path: str) -> str:
    """Extract the '## 基础提示词模板' section from a style markdown file."""
    with open(style_path, "r", encoding="utf-8") as f:
        content = f.read()

    base_prompt_marker = "## 基础提示词模板"
    start_idx = content.find(base_prompt_marker)

    if start_idx == -1:
        print("Warning: '## 基础提示词模板' section not found, using fallback extraction")
        start_idx = content.find("## ")
        end_idx = content.find("## ", start_idx + 3)
        if start_idx == -1 or end_idx == -1:
            return content
        return content[start_idx + 3:end_idx].strip()

    section_start = start_idx + len(base_prompt_marker)
    next_section_idx = content.find("\n## ", section_start)

    if next_section_idx == -1:
        extracted = content[section_start:]
    else:
        extracted = content[section_start:next_section_idx]

    return extracted.strip()


def _style_layout_sidecar_path(style_path: str) -> Path:
    path = Path(style_path)
    if path.suffix:
        return path.with_suffix(".layouts.json")
    return Path(f"{style_path}.layouts.json")


def load_style_layout_profile(
    style_path: str,
    style_template: str = "",
) -> Optional[Dict[str, Any]]:
    """Load optional styles/<id>.layouts.json as a TemplateProfile-compatible dict."""
    sidecar = _style_layout_sidecar_path(style_path)
    if not sidecar.is_file():
        return None

    with open(sidecar, "r", encoding="utf-8") as f:
        profile = json.load(f)
    if not isinstance(profile, dict):
        raise ValueError(f"Style layout sidecar must be a JSON object: {sidecar}")

    layouts = profile.get("layouts")
    if not isinstance(layouts, list) or not layouts:
        raise ValueError(f"Style layout sidecar must contain non-empty layouts: {sidecar}")

    style_summary = str(profile.get("global_style") or "").strip()
    if style_template and style_summary:
        profile["global_style"] = f"{style_template}\n\n【内置 layout bank 风格摘要】\n{style_summary}"
    elif style_template:
        profile["global_style"] = style_template
    else:
        profile["global_style"] = style_summary

    profile.setdefault("version", "2")
    profile.setdefault("source", sidecar.name)
    profile.setdefault("source_hash", "")
    profile.setdefault("theme", {})
    profile.setdefault("style_id", Path(style_path).stem)
    profile["is_style_layout_bank"] = True

    valid_types = {"cover", "agenda", "section", "content", "data", "quote", "closing", "other"}
    for idx, layout in enumerate(layouts):
        if not isinstance(layout, dict):
            raise ValueError(f"Style layout entry #{idx + 1} must be an object: {sidecar}")
        layout.setdefault("id", f"layout-{idx + 1:02d}")
        layout.setdefault("page_index", idx)
        if layout.get("page_type") not in valid_types:
            layout["page_type"] = "content"
        layout.setdefault("summary", "")
        layout["visual_signature"] = str(layout.get("visual_signature") or "").strip()
        capacity = layout.get("content_capacity")
        layout["content_capacity"] = capacity if isinstance(capacity, (dict, list, str)) else {}
        for key in ("best_for", "avoid_for", "variation_tags"):
            val = layout.get(key)
            if isinstance(val, list):
                layout[key] = [str(x).strip() for x in val if str(x).strip()]
            elif isinstance(val, str) and val.strip():
                layout[key] = [val.strip()]
            else:
                layout[key] = []
        layout["external_image_slots"] = (
            layout.get("external_image_slots")
            if isinstance(layout.get("external_image_slots"), list)
            else []
        )
        layout.setdefault("reuse_friendly", layout.get("page_type") != "cover")
        layout.setdefault("reuse_reason", "")
        layout["reference_image"] = layout.get("reference_image") or None
        layout.setdefault("json_schema", {
            "type": "object",
            "properties": {
                "title": {"type": "string", "minLength": 1, "maxLength": 40},
                "body": {"type": "string", "minLength": 0, "maxLength": 600},
            },
            "required": ["title"],
            "additionalProperties": True,
        })

    return profile


def _normalize_template_profile_references(profile: Dict[str, Any], profile_path: str) -> Dict[str, Any]:
    """Resolve relative layout reference_image paths next to template_profile.json."""
    base_dir = Path(profile_path).resolve().parent
    for layout in profile.get("layouts", []) or []:
        ref = layout.get("reference_image")
        if not ref or os.path.isabs(str(ref)):
            continue
        candidate = (base_dir / str(ref)).resolve()
        if candidate.exists():
            layout["reference_image"] = str(candidate)
    return profile


def attach_template_layout_profile(
    slide_spec: Optional[Dict[str, Any]],
    matched_layout: Optional[Dict[str, Any]],
) -> Dict[str, Any]:
    """Attach compact template layout metadata to a slide_spec.

    The full TemplateProfile may include absolute reference image paths and
    schema details. Metadata only needs the layout identity and decision hints
    used for later slot planning, editing, and audit.
    """
    spec = copy.deepcopy(slide_spec) if isinstance(slide_spec, dict) else {}
    if not matched_layout:
        return spec

    profile_keys = [
        "id",
        "page_type",
        "summary",
        "visual_signature",
        "content_capacity",
        "best_for",
        "avoid_for",
        "variation_tags",
        "external_image_slots",
        "reuse_friendly",
        "reuse_reason",
    ]
    compact_profile = {
        key: copy.deepcopy(matched_layout.get(key))
        for key in profile_keys
        if matched_layout.get(key) not in (None, "", [], {})
    }
    spec["template_layout_profile"] = compact_profile

    if not spec.get("layout"):
        layout_id = matched_layout.get("id") or "unknown"
        summary = str(matched_layout.get("summary") or "").strip()
        spec["layout"] = f"模板 {layout_id}：{summary}" if summary else f"模板 {layout_id}"
    return spec


# =============================================================================
# Prompt Generation
# =============================================================================

LANGUAGE_FONT_RULE = """

【强制语言与字体要求】
1. 幻灯片上所有文字必须使用简体中文，严禁出现任何英文单词或句子（产品名称等专有名词可保留英文，其余一律用中文）。
2. 中文字体使用思源黑体（Source Han Sans）或苹方（PingFang SC），字形清晰、笔画规整，严禁使用草书、艺术字或变形字体。
3. 标题字体粗体，正文字体常规，字号对比清晰，确保在演示场景下可读性极高。
"""


def generate_prompt(
    style_template: str,
    page_type: str,
    content_text: str,
    slide_number: int,
    total_slides: int,
    slide_spec: Optional[Dict[str, Any]] = None,
    output_dir: Optional[str] = None,
) -> str:
    """Generate a complete prompt for a single slide.

    When slide_spec is provided (the Agent constructed it from plan + style),
    it produces a precise element-by-element prompt.  Otherwise falls back to
    the original freeform content-based prompt.

    Built-in style prompts provide coarse composition rules; template clone mode
    uses TemplateProfile page types and layout metadata before this fallback path.
    """
    if slide_spec and slide_spec.get("elements"):
        return generate_prompt_from_spec(
            style_template, slide_spec, page_type, slide_number, total_slides,
            output_dir=output_dir,
        )

    is_cover = page_type == "cover" or slide_number == 1
    is_data = page_type == "data" or slide_number == total_slides
    if is_cover:
        label = "封面页（cover）"
        hint = "标题/副标题处理为视觉焦点，按本风格的封面构图规范处理。"
    elif is_data:
        label = "数据页（data）"
        hint = "突出关键数字、对比或结论；按本风格的数据/总结构图规范处理。"
    else:
        label = "内容页（content）"
        hint = "把要点按本风格的内容构图规范结构化呈现，注意层级、对齐、留白。"

    return (
        style_template
        + "\n\n---\n\n"
        + f"现在请生成本组中的【{label}】，{hint}\n"
        + "本页要呈现的内容如下（请按本风格美学重新设计版式，不要原样照搬文本节奏）：\n\n"
        + content_text
        + LANGUAGE_FONT_RULE
    )


# =============================================================================
# External image slot support
# =============================================================================

def _parse_bbox(value: Any) -> Optional[List[float]]:
    """Parse a normalized [x, y, w, h] bbox from slide_spec data."""
    if isinstance(value, (list, tuple)) and len(value) == 4:
        try:
            nums = [float(v) for v in value]
        except (TypeError, ValueError):
            return None
        x, y, w, h = nums
        if w <= 0 or h <= 0:
            return None
        # Clamp to the slide bounds; bad input should degrade, not crash a deck.
        x = max(0.0, min(1.0, x))
        y = max(0.0, min(1.0, y))
        w = max(0.0, min(1.0 - x, w))
        h = max(0.0, min(1.0 - y, h))
        return [x, y, w, h] if w > 0 and h > 0 else None
    return None


def _is_reference_only_image_element(elem: Dict[str, Any]) -> bool:
    etype = str(elem.get("type", "")).lower()
    mode = str(elem.get("render_mode") or elem.get("mode") or "").lower()
    if etype in REFERENCE_ONLY_IMAGE_TYPES:
        return True
    if etype in EXTERNAL_IMAGE_TYPES and mode in REFERENCE_ONLY_MODES:
        return True
    preserve_original = _truthy(elem.get("preserve_original"))
    if etype in EXTERNAL_IMAGE_TYPES and preserve_original is False:
        return True
    return False


def _collect_generation_reference_images(slide_spec: Optional[Dict[str, Any]], output_dir: str) -> List[str]:
    """Return source images used as generation references, not PPT overlays."""
    if not slide_spec:
        return []
    refs: List[str] = []
    for elem in (slide_spec.get("elements") or {}).values():
        if not isinstance(elem, dict) or not _is_reference_only_image_element(elem):
            continue
        source = elem.get("source") or elem.get("path") or elem.get("asset")
        resolved = _resolve_asset_path(source, output_dir)
        if resolved and resolved not in refs:
            refs.append(resolved)
    return refs


def _merge_reference_images(
    existing: Optional[Union[str, List[str]]],
    additions: List[str],
) -> Optional[Union[str, List[str]]]:
    if not additions:
        return existing
    refs: List[str] = []
    if isinstance(existing, (list, tuple)):
        refs.extend(str(p) for p in existing if p)
    elif existing:
        refs.append(str(existing))
    refs.extend(str(p) for p in additions if p)
    deduped: List[str] = []
    for ref in refs:
        if ref not in deduped:
            deduped.append(ref)
    if not deduped:
        return None
    return deduped[0] if len(deduped) == 1 else deduped


def _format_generation_reference_constraint(reference_images: List[str]) -> str:
    if not reference_images:
        return ""
    return (
        "\n\n【真实图参考模式】\n"
        f"本页提供了 {len(reference_images)} 张真实图作为生成参考。用户已允许模型对其进行风格化、融合、重绘或版式再设计，"
        "因此这些图片不会在 PPTX 中作为独立对象原样后贴。"
        "请参考其主体、构图、材质、色彩或场景气质来设计页面，但不要承诺像素级保真。"
        "如果原图中包含文字、数字、表格、医学影像、证据截图、诊断图、工程图或任何需要精确读数/识别的内容，"
        "参考模式不适合保真呈现；这类素材应改用 external_image 独立后贴。"
    )


def _collect_external_image_slots(slide_spec: Optional[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Return normalized external image slot descriptors from a slide_spec.

    Supported element shape:
      {
        "type": "external_image",
        "source": "/path/to/photo.png",
        "position": [0.60, 0.10, 0.32, 0.76],
        "fit": "cover" | "contain",
        "slot": {"padding": 0.012, "fill": "#F7F7F5", "outline": "#1A1A1A"}
      }

    `external_image_placeholder` is also recognized for reference skeletons, but
    final PPTX overlay requires a `source`.
    """
    if not slide_spec:
        return []
    slots: List[Dict[str, Any]] = []
    for elem_id, elem in (slide_spec.get("elements") or {}).items():
        if not isinstance(elem, dict):
            continue
        etype = str(elem.get("type", "")).lower()
        if etype not in EXTERNAL_IMAGE_TYPES:
            continue
        if _is_reference_only_image_element(elem):
            continue
        bbox = _parse_bbox(elem.get("computed_bbox") or elem.get("bbox") or elem.get("position"))
        if not bbox:
            continue
        slot_style = elem.get("slot") if isinstance(elem.get("slot"), dict) else {}
        slots.append({
            "id": str(elem_id),
            "type": etype,
            "bbox": bbox,
            "placement_region": _parse_bbox(elem.get("placement_region") or elem.get("position")),
            "source": elem.get("source") or elem.get("path") or elem.get("asset"),
            "fit": str(elem.get("fit", "contain")).lower(),
            "slot_strategy": str(elem.get("slot_strategy", slot_style.get("strategy", ""))).lower(),
            "tailored_to_asset": bool(elem.get("tailored_to_asset", False)),
            "asset_ratio": elem.get("asset_ratio"),
            "asset_size": elem.get("asset_size"),
            "layout_intent": elem.get("layout_intent"),
            "auto_layout_reason": elem.get("auto_layout_reason"),
            "layout_planning_profile": elem.get("layout_planning_profile"),
            "asset_class": elem.get("asset_class"),
            "detail_level": elem.get("detail_level"),
            "asset_tags": elem.get("asset_tags"),
            "visual_metrics": elem.get("visual_metrics"),
            "trim_whitespace": elem.get("trim_whitespace", slot_style.get("trim_whitespace", "auto")),
            "padding": float(elem.get("padding", slot_style.get("padding", 0.0)) or 0.0),
            "bleed": float(elem.get("bleed", slot_style.get("bleed", 0.0)) or 0.0),
            "mask_bleed": float(elem.get("mask_bleed", slot_style.get("mask_bleed", 0.0)) or 0.0),
            "fill": elem.get("fill", slot_style.get("fill", "#F7F7F5")),
            "mask_fill": elem.get("mask_fill", slot_style.get("mask_fill", slot_style.get("fill", "#F7F7F5"))),
            "outline": elem.get("outline", slot_style.get("outline", "#1A1A1A")),
            "draw_frame": bool(_truthy(elem.get("draw_frame", slot_style.get("draw_frame"))) or False),
            "mask_placeholder": _truthy(elem.get("mask_placeholder", slot_style.get("mask_placeholder", False))) is True,
            "sanitize_background": _truthy(
                elem.get("sanitize_background", slot_style.get("sanitize_background", False))
            ) is True,
            "cleanup_rect": elem.get("cleanup_rect", slot_style.get("cleanup_rect")),
            "outline_width": int(elem.get("outline_width", slot_style.get("outline_width", 0)) or 0),
            "skeleton_canvas_fill": elem.get("skeleton_canvas_fill", slot_style.get("skeleton_canvas_fill", "transparent")),
            "skeleton_fill": elem.get("skeleton_fill", slot_style.get("skeleton_fill", "transparent")),
            "skeleton_outline": elem.get("skeleton_outline", slot_style.get("skeleton_outline", "#000000")),
            "skeleton_outline_width": int(
                elem.get(
                    "skeleton_outline_width",
                    slot_style.get("skeleton_outline_width", slot_style.get("reference_outline_width", 2)),
                ) or 0
            ),
            "skeleton_shape": str(
                elem.get("skeleton_shape", slot_style.get("skeleton_shape", "corners"))
            ).lower(),
            "skeleton_ticks": _truthy(elem.get("skeleton_ticks", slot_style.get("skeleton_ticks", False))) is True,
        })
    return slots


def _slots_with_real_sources(slots: List[Dict[str, Any]], output_dir: str) -> List[Dict[str, Any]]:
    """Keep only slots that point to a real local image asset."""
    return [
        slot for slot in slots
        if _resolve_asset_path(slot.get("source"), output_dir)
    ]


def _hex_to_rgb(value: Any, default: Tuple[int, int, int] = (247, 247, 245)) -> Tuple[int, int, int]:
    if not isinstance(value, str):
        return default
    text = value.strip()
    if text.lower() in {"none", "transparent"}:
        return default
    if text.startswith("#"):
        text = text[1:]
    if len(text) == 3:
        text = "".join(ch * 2 for ch in text)
    if len(text) != 6:
        return default
    try:
        return tuple(int(text[i:i + 2], 16) for i in (0, 2, 4))  # type: ignore[return-value]
    except ValueError:
        return default


def _is_transparent_color(value: Any) -> bool:
    return isinstance(value, str) and value.strip().lower() in {"none", "transparent", "no-fill", "nofill"}


def _resolve_asset_path(source: Any, output_dir: str) -> Optional[str]:
    if not source or not isinstance(source, str):
        return None
    p = Path(source).expanduser()
    if not p.is_absolute():
        p = Path(output_dir) / p
    return str(p.resolve()) if p.is_file() else None


def _norm_rect_to_pixels(rect: List[float], width: int = SLIDE_TRACE_WIDTH, height: int = SLIDE_TRACE_HEIGHT) -> List[int]:
    x, y, w, h = rect
    return [
        int(round(x * width)),
        int(round(y * height)),
        int(round(w * width)),
        int(round(h * height)),
    ]


def _slot_inner_rect(slot: Dict[str, Any]) -> List[float]:
    x, y, w, h = slot["bbox"]
    pad = max(0.0, float(slot.get("padding") or 0.0))
    bleed = max(0.0, float(slot.get("bleed") or 0.0))
    inner_x = max(0.0, x + pad - bleed)
    inner_y = max(0.0, y + pad - bleed)
    inner_w = max(0.001, min(1.0 - inner_x, w - 2 * pad + 2 * bleed))
    inner_h = max(0.001, min(1.0 - inner_y, h - 2 * pad + 2 * bleed))
    return [inner_x, inner_y, inner_w, inner_h]


def _compute_final_image_rect_px(slot: Dict[str, Any], asset_path: str, width: int, height: int) -> List[int]:
    """Return the actual picture rectangle used by contain/cover placement."""
    inner = _slot_inner_rect(slot)
    box_left, box_top, box_w, box_h = _norm_rect_to_pixels(inner, width, height)
    if slot.get("fit") == "cover":
        return [box_left, box_top, box_w, box_h]

    try:
        from PIL import Image
        with Image.open(asset_path) as im:
            ratio = im.width / im.height
    except Exception:
        ratio = box_w / box_h if box_h else 1.0

    box_ratio = box_w / box_h if box_h else ratio
    if ratio >= box_ratio:
        draw_w = box_w
        draw_h = int(draw_w / ratio)
        draw_left = box_left
        draw_top = box_top + int((box_h - draw_h) / 2)
    else:
        draw_h = box_h
        draw_w = int(draw_h * ratio)
        draw_left = box_left + int((box_w - draw_w) / 2)
        draw_top = box_top
    return [draw_left, draw_top, draw_w, draw_h]


def _paste_asset_into_rect(canvas: Any, asset_path: str, rect_px: List[int], fit: str = "contain") -> None:
    from PIL import Image

    left, top, draw_w, draw_h = rect_px
    if draw_w <= 0 or draw_h <= 0:
        return
    asset = Image.open(asset_path).convert("RGB")
    if fit == "cover":
        target_ratio = draw_w / draw_h
        ratio = asset.width / asset.height
        if ratio > target_ratio:
            new_w = int(asset.height * target_ratio)
            x0 = max(0, (asset.width - new_w) // 2)
            asset = asset.crop((x0, 0, x0 + new_w, asset.height))
        elif ratio < target_ratio:
            new_h = int(asset.width / target_ratio)
            y0 = max(0, (asset.height - new_h) // 2)
            asset = asset.crop((0, y0, asset.width, y0 + new_h))
    canvas.paste(asset.resize((draw_w, draw_h)), (left, top))


def _trace_dir(output_dir: str, slide_number: int) -> Path:
    return Path(output_dir) / "external_image_trace" / f"slide-{slide_number:02d}"


def _write_json(path: Path, data: Dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


def _load_trace_manifest(path: Path) -> Dict[str, Any]:
    if path.is_file():
        try:
            return json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            return {}
    return {}


def _truthy(value: Any) -> Optional[bool]:
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return bool(value)
    if isinstance(value, str):
        text = value.strip().lower()
        if text in {"1", "true", "yes", "y", "on"}:
            return True
        if text in {"0", "false", "no", "n", "off"}:
            return False
    return None


def _asset_size_and_ratio(source: Any, output_dir: str) -> Optional[Tuple[int, int, float]]:
    path = _resolve_asset_path(source, output_dir)
    if not path:
        return None
    try:
        from PIL import Image
        with Image.open(path) as im:
            if im.width <= 0 or im.height <= 0:
                return None
            return im.width, im.height, im.width / im.height
    except Exception as e:
        print(f"(!)  无法读取外部图片尺寸，跳过量体裁衣: {path} ({e})")
        return None


def _fit_bbox_to_asset_ratio(
    region: List[float],
    asset_ratio: float,
    padding: float = 0.0,
    anchor: str = "center",
) -> List[float]:
    """Fit a slot inside a normalized slide region while preserving asset ratio.

    The existing overlay code treats `padding` as normalized x/y in slide
    coordinates, so the inner normalized ratio must be adjusted by the 16:9
    slide aspect to match the image's actual pixel ratio.
    """
    slide_aspect = 16.0 / 9.0
    x, y, w, h = region
    pad = max(0.0, min(float(padding or 0.0), min(w, h) / 2.0 - 0.0005))
    inner_max_w = max(0.001, w - 2 * pad)
    inner_max_h = max(0.001, h - 2 * pad)
    target_norm_ratio = max(0.001, asset_ratio / slide_aspect)

    if inner_max_w / inner_max_h >= target_norm_ratio:
        inner_h = inner_max_h
        inner_w = inner_h * target_norm_ratio
    else:
        inner_w = inner_max_w
        inner_h = inner_w / target_norm_ratio

    outer_w = min(w, inner_w + 2 * pad)
    outer_h = min(h, inner_h + 2 * pad)

    anchor_text = (anchor or "center").lower().replace("_", "-")
    if "left" in anchor_text:
        out_x = x
    elif "right" in anchor_text:
        out_x = x + w - outer_w
    else:
        out_x = x + (w - outer_w) / 2

    if "top" in anchor_text:
        out_y = y
    elif "bottom" in anchor_text:
        out_y = y + h - outer_h
    else:
        out_y = y + (h - outer_h) / 2

    return [
        max(0.0, min(1.0, out_x)),
        max(0.0, min(1.0, out_y)),
        max(0.001, min(1.0 - out_x, outer_w)),
        max(0.001, min(1.0 - out_y, outer_h)),
    ]


def _external_element_items(slide_spec: Dict[str, Any]) -> List[Tuple[str, Dict[str, Any]]]:
    items: List[Tuple[str, Dict[str, Any]]] = []
    for elem_id, elem in (slide_spec.get("elements") or {}).items():
        if not isinstance(elem, dict):
            continue
        etype = str(elem.get("type", "")).lower()
        if etype in EXTERNAL_IMAGE_TYPES and not _is_reference_only_image_element(elem):
            items.append((str(elem_id), elem))
    return items


def _slide_text_load(slide_spec: Dict[str, Any]) -> int:
    chars = 0
    for elem in (slide_spec.get("elements") or {}).values():
        if not isinstance(elem, dict):
            continue
        etype = str(elem.get("type", "")).lower()
        if etype in EXTERNAL_IMAGE_TYPES or etype in REFERENCE_ONLY_IMAGE_TYPES:
            continue
        chars += len(_element_text_content(elem, include_description=True))
    layout_text = str(slide_spec.get("layout") or "")
    return chars + min(len(layout_text), 80)


def _slide_text_profile(slide_spec: Dict[str, Any]) -> Dict[str, Any]:
    title_chars = 0
    body_chars = 0
    text_blocks = 0
    for elem in (slide_spec.get("elements") or {}).values():
        if not isinstance(elem, dict):
            continue
        etype = str(elem.get("type", "")).lower()
        if etype in EXTERNAL_IMAGE_TYPES or etype in REFERENCE_ONLY_IMAGE_TYPES:
            continue
        length = len(_element_text_content(elem, include_description=True))
        if not length:
            continue
        text_blocks += 1
        if etype in {"title", "headline", "heading", "subtitle"}:
            title_chars += length
        else:
            body_chars += length
    layout_chars = min(len(str(slide_spec.get("layout") or "")), 80)
    text_load = title_chars + body_chars + layout_chars
    if text_load >= 180:
        density = "heavy"
    elif text_load >= 90:
        density = "medium"
    else:
        density = "light"
    return {
        "title_chars": title_chars,
        "body_chars": body_chars,
        "layout_chars": layout_chars,
        "text_blocks": text_blocks,
        "text_load": text_load,
        "density": density,
    }


def _is_auto_external_layout(elem: Dict[str, Any]) -> bool:
    intent = str(elem.get("layout_intent", elem.get("auto_layout", ""))).strip().lower()
    position = elem.get("position")
    return (
        intent in {"auto", "auto-layout", "autolayout", "smart", "smart-layout"}
        or (isinstance(position, str) and position.strip().lower() in {"auto", "auto-layout", "smart"})
        or position is None
    )


def _analyze_external_asset_visual(source: Any, output_dir: str) -> Dict[str, Any]:
    """Heuristically classify a real image before layout planning.

    This is deliberately local and deterministic. Vision review remains the
    final guard, but these tags let the first generation start with a better
    slot size instead of learning only after a failed overlay review.
    """
    path = _resolve_asset_path(source, output_dir)
    if not path:
        return {}
    try:
        from PIL import Image, ImageStat

        with Image.open(path) as im:
            original_w, original_h = im.size
            rgb = im.convert("RGB")
            rgb.thumbnail((256, 256))
            stat = ImageStat.Stat(rgb)
            mean = stat.mean
            stddev = stat.stddev
            pixels = list(rgb.getdata())
            total = max(1, len(pixels))
            white_ratio = sum(1 for r, g, b in pixels if r > 238 and g > 238 and b > 238) / total
            dark_ratio = sum(1 for r, g, b in pixels if r < 80 and g < 80 and b < 80) / total
            color_std = sum(stddev) / 3.0

            gray = rgb.convert("L")
            w, h = gray.size
            edge_hits = 0
            edge_total = 0
            if w > 1 and h > 1:
                pix = gray.load()
                step = max(1, min(w, h) // 96)
                for y in range(0, h - step, step):
                    for x in range(0, w - step, step):
                        gx = abs(int(pix[x + step, y]) - int(pix[x, y]))
                        gy = abs(int(pix[x, y + step]) - int(pix[x, y]))
                        if gx + gy > 42:
                            edge_hits += 1
                        edge_total += 1
            edge_density = edge_hits / max(1, edge_total)
    except Exception:
        return {}

    ratio = original_w / original_h if original_h else 1.0
    tags: List[str] = []
    if ratio >= 2.0:
        tags.append("wide")
    elif ratio <= 0.8:
        tags.append("portrait")
    else:
        tags.append("standard-ratio")
    if white_ratio > 0.45:
        tags.append("white-background")
    if edge_density > 0.16:
        tags.append("dense-lines")
    elif edge_density > 0.09:
        tags.append("line-art")
    if dark_ratio > 0.08:
        tags.append("text-like")
    if color_std > 58 and white_ratio < 0.35:
        tags.append("photo-like")

    if "photo-like" in tags and edge_density < 0.12:
        asset_class = "photo"
        detail_level = "normal"
    elif ratio >= 2.0 and (white_ratio > 0.35 or dark_ratio > 0.06):
        asset_class = "wide-document"
        detail_level = "high"
    elif white_ratio > 0.65 and edge_density > 0.04 and color_std < 35:
        asset_class = "dense-diagram"
        detail_level = "critical"
    elif white_ratio > 0.45 and edge_density > 0.13:
        asset_class = "dense-diagram"
        detail_level = "critical"
    elif white_ratio > 0.35 and (edge_density > 0.08 or dark_ratio > 0.08):
        asset_class = "chart-or-document"
        detail_level = "high"
    else:
        asset_class = "general-image"
        detail_level = "normal"

    return {
        "asset_class": asset_class,
        "detail_level": detail_level,
        "asset_tags": tags,
        "visual_metrics": {
            "white_ratio": round(white_ratio, 4),
            "dark_ratio": round(dark_ratio, 4),
            "edge_density": round(edge_density, 4),
            "color_std": round(color_std, 2),
            "mean_rgb": [round(x, 2) for x in mean],
        },
    }


def _external_asset_detail_level(slide_spec: Dict[str, Any], elem_id: str, elem: Dict[str, Any]) -> str:
    """Classify whether a real asset needs enough size to be read, not just seen."""
    explicit = str(
        elem.get("detail_level")
        or elem.get("readability")
        or elem.get("readability_priority")
        or ""
    ).strip().lower()
    if explicit in {"low", "photo", "decorative", "ambient"}:
        return "low"

    asset_class = str(elem.get("asset_class") or "").lower()
    if explicit in {"critical", "very-high", "dense", "dense-diagram"}:
        return "critical"
    if asset_class in {"dense-diagram", "dense-chart", "architecture-diagram"}:
        return "critical"
    if explicit in {"high", "detail", "detailed", "readable", "legible", "text-heavy", "diagram"}:
        return "high"
    if asset_class in {"wide-document", "chart-or-document", "screenshot-document"}:
        return "high"

    context = " ".join(
        str(x)
        for x in [
            elem_id,
            elem.get("label", ""),
            elem.get("caption", ""),
            elem.get("description", ""),
            slide_spec.get("layout", ""),
        ]
    ).lower()
    high_terms = {
        "paper", "table", "proof", "architecture", "diagram", "chart", "curve",
        "graph", "screenshot", "code", "formula", "equation", "lean", "result",
        "论文", "表格", "证明", "架构", "流程", "图表", "曲线", "截图", "代码",
        "公式", "定理", "结果", "系统图",
    }
    if any(term in context for term in high_terms):
        return "high"
    return "normal"


def _region_area(region: List[float]) -> float:
    return max(0.0, float(region[2])) * max(0.0, float(region[3]))


def _region_fit_efficiency(region: List[float], ratio: float) -> float:
    """How much of a candidate slot a contain-fitted asset can actually use."""
    if ratio <= 0:
        return 0.0
    x, y, w, h = region
    px_w = w * SLIDE_TRACE_WIDTH
    px_h = h * SLIDE_TRACE_HEIGHT
    if px_w <= 1 or px_h <= 1:
        return 0.0
    slot_ratio = px_w / px_h
    if slot_ratio >= ratio:
        used_w = px_h * ratio
        used_h = px_h
    else:
        used_w = px_w
        used_h = px_w / ratio
    return max(0.0, min(1.0, (used_w * used_h) / (px_w * px_h)))


def _normalize_candidate_region(region: Any) -> Optional[List[float]]:
    parsed = _parse_bbox(region)
    if not parsed:
        return None
    x, y, w, h = parsed
    if w < 0.12 or h < 0.12:
        return None
    return [x, y, w, h]


def _template_slot_candidates(slide_spec: Dict[str, Any]) -> List[Tuple[List[float], str, float]]:
    """Return template/style-derived image zones before generic fallback rules.

    The template analyzer can provide explicit normalized media slots. Older
    profiles and built-in styles may only have a layout summary, so we also
    infer a few candidate zones from common composition language.
    """
    candidates: List[Tuple[List[float], str, float]] = []
    profile = slide_spec.get("template_layout_profile")
    summary = ""
    if isinstance(profile, dict):
        summary = " ".join(
            str(x)
            for x in [
                profile.get("id", ""),
                profile.get("page_type", ""),
                profile.get("summary", ""),
            ]
        )
        for idx, slot in enumerate(profile.get("external_image_slots") or []):
            if not isinstance(slot, dict):
                continue
            region = _normalize_candidate_region(slot.get("bbox"))
            if not region:
                continue
            try:
                priority = float(slot.get("priority", idx + 1) or (idx + 1))
            except (TypeError, ValueError):
                priority = idx + 1
            candidates.append((
                region,
                f"template explicit slot {slot.get('id') or idx + 1}: {slot.get('purpose', '')}".strip(),
                max(0.70, 1.25 - 0.12 * (priority - 1)),
            ))

    layout_text = " ".join(
        str(x)
        for x in [
            summary,
            slide_spec.get("layout", ""),
            slide_spec.get("page_type", ""),
        ]
    ).lower()

    def add(region: List[float], reason: str, weight: float) -> None:
        parsed = _normalize_candidate_region(region)
        if parsed:
            candidates.append((parsed, reason, weight))

    # Explicit directional patterns. Phrases like "左文右图" mean the image zone
    # is on the right; "右文左图" means the image zone is on the left.
    if any(t in layout_text for t in [
        "左文右图", "右侧为图", "右侧为图片", "右侧为照片", "图像在右", "图片在右",
        "left text right image", "text left image right", "right image", "image on the right",
        "photo on the right", "right photo", "photo right", "right photograph", "right visual",
        "right landscape photo", "right preserved landscape photo",
        "media on the right",
    ]):
        add([0.52, 0.10, 0.42, 0.78], "template inferred right image column", 1.10)
    if any(t in layout_text for t in [
        "右文左图", "左侧为图", "左侧为图片", "左侧为照片", "图像在左", "图片在左",
        "right text left image", "text right image left", "left image", "image on the left",
        "photo on the left", "left photo", "photo left", "left photograph", "left visual",
        "left landscape photo", "left preserved landscape photo",
        "media on the left",
    ]):
        add([0.06, 0.10, 0.42, 0.78], "template inferred left image column", 1.10)
    if any(t in layout_text for t in ["图片区在左", "左侧图片", "左侧照片", "左侧主视觉", "left visual", "left media"]):
        add([0.06, 0.10, 0.44, 0.78], "template inferred left visual region", 1.08)
    if any(t in layout_text for t in ["图片区在右", "右侧图片", "右侧照片", "右侧主视觉", "right visual", "right media"]):
        add([0.50, 0.10, 0.44, 0.78], "template inferred right visual region", 1.08)
    if any(t in layout_text for t in ["底部图片", "底部图表", "下方图片", "bottom image", "bottom chart", "lower image"]):
        add([0.08, 0.56, 0.84, 0.36], "template inferred bottom media band", 1.04)
    if any(t in layout_text for t in ["顶部图片", "顶部图表", "上方图片", "top image", "top chart", "upper image"]):
        add([0.08, 0.10, 0.84, 0.36], "template inferred top media band", 1.02)
    if any(t in layout_text for t in ["中央主视觉", "居中主视觉", "中心图片", "central visual", "center visual", "centered image"]):
        add([0.22, 0.14, 0.56, 0.68], "template inferred central visual region", 1.02)
    if any(t in layout_text for t in ["分屏", "split screen", "split-screen"]):
        add([0.52, 0.08, 0.42, 0.84], "template inferred split-screen right panel", 0.96)
        add([0.06, 0.08, 0.42, 0.84], "template inferred split-screen left panel", 0.94)

    # De-duplicate near-identical candidates while keeping the highest weight.
    deduped: List[Tuple[List[float], str, float]] = []
    for region, reason, weight in candidates:
        for idx, (existing, existing_reason, existing_weight) in enumerate(deduped):
            if _bbox_delta(existing, region) < 0.08:
                if weight > existing_weight:
                    deduped[idx] = (region, reason, weight)
                break
        else:
            deduped.append((region, reason, weight))
    return deduped


def _choose_template_aware_region(
    slide_spec: Dict[str, Any],
    elem_id: str,
    elem: Dict[str, Any],
    asset_ratio: Optional[float],
    text_profile: Dict[str, Any],
    total: int,
    index: int,
    detail_level: str,
    planning_bits: str,
) -> Optional[Tuple[List[float], str]]:
    candidates = _template_slot_candidates(slide_spec)
    if not candidates:
        return None
    ratio = asset_ratio or 1.6
    text_load = int(text_profile["text_load"])
    density = str(text_profile["density"])
    min_area = 0.18
    if detail_level == "high":
        min_area = 0.26
    elif detail_level == "critical":
        min_area = 0.36
    if total >= 2:
        min_area *= 0.62

    best: Optional[Tuple[float, List[float], str]] = None
    for region, reason, weight in candidates:
        x, y, w, h = region
        area = _region_area(region)
        fit = _region_fit_efficiency(region, ratio)
        if detail_level in {"high", "critical"} and fit < 0.52:
            # A template may expose a very wide media band, but a standard-ratio
            # chart/table would render as a small thumbnail inside it. In that
            # case preserving the template slot is worse than choosing a more
            # balanced readable region.
            continue
        score = weight + min(area / max(min_area, 0.01), 1.35) * 0.55 + fit * 0.35
        if detail_level in {"high", "critical"} and area < min_area:
            score -= (min_area - area) * 2.0
        if density == "heavy" and y < 0.18 and h > 0.65 and 0.18 < x < 0.45:
            score -= 0.18
        if total == 1 and detail_level in {"high", "critical"} and area >= min_area:
            score += 0.12
        if best is None or score > best[0]:
            best = (score, region, f"{reason}; score={score:.2f}; {planning_bits}")

    if not best:
        return None
    score, region, reason = best
    # Do not force a tiny template thumbnail slot for dense assets. In that
    # case, let the generic readability fallback choose a larger area.
    if detail_level in {"high", "critical"} and _region_area(region) < min_area * 0.82:
        return None
    return region, f"template-aware slot selected: {reason}"


def _auto_external_image_region(
    slide_spec: Dict[str, Any],
    elem_id: str,
    elem: Dict[str, Any],
    asset_ratio: Optional[float],
    index: int,
    total: int,
) -> Tuple[List[float], str]:
    """Choose an allowed region for a real image before generating the slide."""
    ratio = asset_ratio or 1.6
    text_profile = _slide_text_profile(slide_spec)
    text_load = int(text_profile["text_load"])
    title_chars = int(text_profile["title_chars"])
    density = str(text_profile["density"])
    page_type = str(slide_spec.get("page_type", "") or "").lower()
    detail_level = _external_asset_detail_level(slide_spec, elem_id, elem)
    planning_context = " ".join(
        str(x)
        for x in [
            elem_id,
            elem.get("label", ""),
            elem.get("caption", ""),
            elem.get("description", ""),
            slide_spec.get("layout", ""),
        ]
    ).lower()
    structural_terms = {
        "architecture", "diagram", "flow", "workflow", "system", "agent",
        "transformer", "pipeline", "loop", "架构", "流程", "系统", "闭环",
        "模块", "链路",
    }
    is_structural_diagram = any(term in planning_context for term in structural_terms)
    planning_bits = (
        f"text_density={density}; title_chars={title_chars}; "
        f"text_load={text_load}; image_count={total}; ratio={ratio:.2f}; detail={detail_level}"
    )

    template_region = _choose_template_aware_region(
        slide_spec, elem_id, elem, asset_ratio, text_profile, total, index, detail_level, planning_bits
    )
    if template_region:
        return template_region

    if total >= 3:
        col_w = 0.78 / min(total, 3)
        col = index % 3
        row = index // 3
        y = 0.62 if density in {"medium", "heavy"} and row == 0 else (0.58 if row == 0 else 0.34)
        return [0.11 + col * col_w, y, col_w - 0.025, 0.26], (
            f"auto multi-image grid after text-first planning; {planning_bits}"
        )

    if total == 2:
        if detail_level in {"high", "critical"}:
            if ratio < 0.8:
                return [0.48 + index * 0.24, 0.10, 0.22, 0.78], (
                    f"auto two readable portrait rails; {planning_bits}"
                )
            return [0.08 + index * 0.44, 0.54, 0.40, 0.36], (
                f"auto two readable bottom panels; {planning_bits}"
            )
        if ratio < 0.8:
            return [0.08 + index * 0.28, 0.18, 0.24, 0.64], (
                f"auto two portrait side-by-side; {planning_bits}"
            )
        if density == "light":
            return [0.10 + index * 0.42, 0.52, 0.38, 0.34], (
                f"auto two-image lower gallery with open title area; {planning_bits}"
            )
        return [0.56, 0.15 + index * 0.35, 0.36, 0.29], (
            f"auto two-image right column; {planning_bits}"
        )

    if detail_level == "critical":
        if ratio >= 2.0:
            return [0.06, 0.48, 0.88, 0.44], (
                f"auto near-full bottom panel for critical wide asset; {planning_bits}"
            )
        if ratio < 0.8:
            return [0.52, 0.08, 0.42, 0.84], (
                f"auto near-full right rail for critical portrait asset; {planning_bits}"
            )
        if not is_structural_diagram and text_load >= 80:
            return [0.34, 0.06, 0.62, 0.88], (
                f"auto enlarged right chart/document panel with compact protected text rail; {planning_bits}"
            )
        return [0.30, 0.08, 0.66, 0.84], (
            f"auto right panel with protected text rail for critical detailed asset; {planning_bits}"
        )

    if detail_level == "high":
        if ratio >= 2.0:
            return [0.08, 0.52, 0.84, 0.40], (
                f"auto large bottom band for detailed wide asset; {planning_bits}"
            )
        if ratio < 0.8:
            return [0.58, 0.10, 0.36, 0.78], (
                f"auto large right rail for detailed portrait asset; {planning_bits}"
            )
        return [0.46, 0.10, 0.48, 0.80], (
            f"auto large right panel for detailed asset; {planning_bits}"
        )

    if text_load >= 180:
        if ratio >= 1.25:
            return [0.56, 0.20, 0.36, 0.50], f"auto right landscape block for text-heavy slide; {planning_bits}"
        if ratio < 0.8:
            return [0.68, 0.17, 0.24, 0.64], f"auto narrow right rail for text-heavy portrait; {planning_bits}"
        return [0.64, 0.22, 0.28, 0.42], f"auto compact right slot for text-heavy slide; {planning_bits}"

    if text_load >= 90:
        if ratio < 0.8:
            return [0.66, 0.14, 0.26, 0.68], f"auto right portrait rail for medium text; {planning_bits}"
        return [0.56, 0.18, 0.36, 0.52], f"auto right image block for medium text; {planning_bits}"

    if page_type == "cover" and ratio >= 1.1:
        return [0.52, 0.18, 0.40, 0.56], f"auto cover hero slot; {planning_bits}"
    if ratio < 0.8:
        return [0.66, 0.14, 0.26, 0.70], f"auto right portrait feature rail; {planning_bits}"
    return [0.54, 0.18, 0.38, 0.54], f"auto standard right slot; {planning_bits}"


def prepare_external_image_slots(slide_spec: Optional[Dict[str, Any]], output_dir: str) -> Optional[Dict[str, Any]]:
    """Compute final external image bboxes before prompt/skeleton/PPTX.

    For `tailor_to_asset` / `slot_strategy: fit-within`, `position` is treated
    as the allowed region. The function writes a code-owned `computed_bbox`
    based on the real asset's aspect ratio, and later stages all consume that
    same bbox. This keeps the skeleton reference and PPTX overlay in sync.
    """
    if not slide_spec or not slide_spec.get("elements"):
        return slide_spec

    prepared = copy.deepcopy(slide_spec)
    external_items = _external_element_items(prepared)
    total_external = len(external_items)
    for idx, (_elem_id, elem) in enumerate(external_items):
        if not isinstance(elem, dict):
            continue
        etype = str(elem.get("type", "")).lower()
        if etype not in EXTERNAL_IMAGE_TYPES or _is_reference_only_image_element(elem):
            continue

        slot_style = elem.get("slot") if isinstance(elem.get("slot"), dict) else {}
        source = elem.get("source") or elem.get("path") or elem.get("asset")
        size_ratio = _asset_size_and_ratio(source, output_dir)
        if size_ratio:
            asset_w, asset_h, asset_ratio = size_ratio
            elem["asset_size"] = [asset_w, asset_h]
            elem["asset_ratio"] = round(asset_ratio, 6)
        else:
            asset_ratio = None
        visual_analysis = _analyze_external_asset_visual(source, output_dir)
        if visual_analysis:
            elem.setdefault("asset_class", visual_analysis.get("asset_class"))
            elem.setdefault("asset_tags", visual_analysis.get("asset_tags"))
            elem["visual_metrics"] = visual_analysis.get("visual_metrics")
            if not any(elem.get(k) for k in ("detail_level", "readability", "readability_priority")):
                elem["detail_level"] = visual_analysis.get("detail_level")

        if _is_auto_external_layout(elem):
            region, reason = _auto_external_image_region(prepared, _elem_id, elem, asset_ratio, idx, total_external)
            elem["position"] = region
            elem["auto_layout_reason"] = reason
            elem["layout_intent"] = "auto"
            elem["layout_planning_profile"] = {
                **_slide_text_profile(prepared),
                "image_count": total_external,
                "image_index": idx,
                "asset_ratio": round(asset_ratio, 6) if asset_ratio else None,
                "detail_level": _external_asset_detail_level(prepared, _elem_id, elem),
                "asset_class": elem.get("asset_class"),
            }
        else:
            region = _parse_bbox(elem.get("position") or elem.get("bbox"))
        if not region:
            continue

        strategy = str(elem.get("slot_strategy", slot_style.get("strategy", ""))).lower()
        tailor_flag = _truthy(elem.get("tailor_to_asset", slot_style.get("tailor_to_asset")))
        fit = str(elem.get("fit", "contain")).lower()

        # Back-compatible rule: explicit bbox remains exact by default. A
        # position-only contain image is usually an allowed region, so tailor it.
        should_tailor = (
            asset_ratio is not None
            and strategy not in {"exact", "fixed", "cover-region"}
            and (
                tailor_flag is True
                or strategy in {"fit-within", "contain-within", "tailored-contain", "asset-fit"}
                or (tailor_flag is None and "bbox" not in elem and fit == "contain")
            )
        )

        elem["placement_region"] = region
        if should_tailor:
            padding = float(elem.get("padding", slot_style.get("padding", 0.0)) or 0.0)
            anchor = str(elem.get("anchor", slot_style.get("anchor", "center")))
            elem["computed_bbox"] = _fit_bbox_to_asset_ratio(region, asset_ratio, padding, anchor)
            elem["tailored_to_asset"] = True
            elem["slot_strategy"] = strategy or "fit-within"
        else:
            elem["computed_bbox"] = _parse_bbox(elem.get("bbox") or elem.get("position"))
            elem["tailored_to_asset"] = False
            if strategy:
                elem["slot_strategy"] = strategy

    return prepared


def create_asset_reference_skeleton(
    output_dir: str,
    slide_number: int,
    slide_spec: Optional[Dict[str, Any]],
) -> Optional[str]:
    """Create a transparent skeleton reference image marking external image slots.

    This reference gives gpt-image-2 a concrete visual constraint for reserved
    space. The final geometry still comes from slide_spec bbox values.
    """
    slide_spec = prepare_external_image_slots(slide_spec, output_dir)
    slots = _slots_with_real_sources(_collect_external_image_slots(slide_spec), output_dir)
    if not slots:
        return None
    try:
        from PIL import Image, ImageDraw
    except ImportError:
        print("(!)  缺 Pillow，跳过外部图片槽位参考图生成（pip install pillow）")
        return None

    width, height = SLIDE_TRACE_WIDTH, SLIDE_TRACE_HEIGHT
    canvas_fill = slots[0].get("skeleton_canvas_fill", "transparent")
    if _is_transparent_color(canvas_fill):
        img = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    else:
        img = Image.new("RGB", (width, height), _hex_to_rgb(canvas_fill, (247, 247, 245)))
    draw = ImageDraw.Draw(img)
    real_on_blank = Image.new("RGB", (width, height), (247, 247, 245))
    trace_slots: List[Dict[str, Any]] = []
    for slot in slots:
        asset_path = _resolve_asset_path(slot.get("source"), output_dir)
        display_asset_path = _display_asset_for_slot(slot, asset_path, output_dir, slide_number) if asset_path else None
        inner_rect = _slot_inner_rect(slot)
        final_rect_px = None
        if display_asset_path:
            final_rect_px = _compute_final_image_rect_px(slot, display_asset_path, width, height)
        else:
            final_rect_px = _norm_rect_to_pixels(inner_rect, width, height)
        rect = [
            final_rect_px[0],
            final_rect_px[1],
            final_rect_px[0] + final_rect_px[2],
            final_rect_px[1] + final_rect_px[3],
        ]
        fill = None if _is_transparent_color(slot.get("skeleton_fill")) else _hex_to_rgb(slot.get("skeleton_fill"), (244, 240, 235))
        outline_width = int(slot.get("skeleton_outline_width") or 0)
        outline = _hex_to_rgb(slot.get("skeleton_outline"), (216, 211, 200))
        skeleton_shape = str(slot.get("skeleton_shape") or "corners").lower()
        if skeleton_shape in {"outline", "rectangle", "rect"} and outline_width > 0:
            draw.rectangle(rect, fill=fill, outline=outline, width=outline_width)
        elif skeleton_shape in {"fill", "filled"}:
            if fill is not None:
                draw.rectangle(rect, fill=fill)
        elif skeleton_shape not in {"none", "off", "false"} and outline_width > 0:
            # Closed rectangles are often interpreted by image models as real
            # photo frames. Corner marks preserve the geometry while making the
            # guide less likely to become a visible placeholder in the output.
            x1, y1, x2, y2 = rect
            mark = max(14, min(x2 - x1, y2 - y1) // 10)
            corners = (
                ((x1, y1), (x1 + mark, y1), (x1, y1 + mark)),
                ((x2, y1), (x2 - mark, y1), (x2, y1 + mark)),
                ((x1, y2), (x1 + mark, y2), (x1, y2 - mark)),
                ((x2, y2), (x2 - mark, y2), (x2, y2 - mark)),
            )
            for corner, horizontal, vertical in corners:
                draw.line((corner, horizontal), fill=outline, width=outline_width)
                draw.line((corner, vertical), fill=outline, width=outline_width)
        if slot.get("skeleton_ticks"):
            # Optional corner ticks are useful for debugging, but should stay
            # off by default because image models may copy them into the slide.
            line_w = max(1, outline_width)
            tick = max(10, min(rect[2] - rect[0], rect[3] - rect[1]) // 20)
            for cx, cy in ((rect[0], rect[1]), (rect[2], rect[1]), (rect[0], rect[3]), (rect[2], rect[3])):
                draw.line((cx - tick, cy, cx + tick, cy), fill=outline, width=line_w)
                draw.line((cx, cy - tick, cx, cy + tick), fill=outline, width=line_w)

        if display_asset_path:
            _paste_asset_into_rect(real_on_blank, display_asset_path, final_rect_px, fit=slot.get("fit", "contain"))

        trace_slots.append({
            "id": slot.get("id"),
            "source": asset_path,
            "display_source": display_asset_path,
            "fit": slot.get("fit"),
            "slot_bbox_norm": slot.get("bbox"),
            "slot_bbox_px": _norm_rect_to_pixels(slot.get("bbox"), width, height),
            "inner_rect_norm": inner_rect,
            "inner_rect_px": _norm_rect_to_pixels(inner_rect, width, height),
            "reference_rect_px": final_rect_px,
            "final_image_rect_px": final_rect_px,
            "padding": slot.get("padding"),
            "bleed": slot.get("bleed"),
            "skeleton_canvas_fill": slot.get("skeleton_canvas_fill"),
            "skeleton_shape": skeleton_shape,
            "layout_intent": slot.get("layout_intent"),
            "auto_layout_reason": slot.get("auto_layout_reason"),
            "tailored_to_asset": slot.get("tailored_to_asset"),
            "asset_size": slot.get("asset_size"),
            "asset_ratio": slot.get("asset_ratio"),
        })

    ref_dir = Path(output_dir) / "references"
    ref_dir.mkdir(parents=True, exist_ok=True)
    ref_path = ref_dir / f"slide-{slide_number:02d}-asset-skeleton.png"
    img.save(ref_path)

    trace = _trace_dir(output_dir, slide_number)
    trace.mkdir(parents=True, exist_ok=True)
    step1 = trace / "step1-real-on-blank.png"
    step2 = trace / "step2-reference-outline-blank.png"
    real_on_blank.save(step1)
    img.save(step2)
    manifest_path = trace / "manifest.json"
    manifest = _load_trace_manifest(manifest_path)
    manifest.update({
        "slide_number": slide_number,
        "canvas_size": [width, height],
        "step1_real_on_blank": str(step1),
        "step2_reference_outline_blank": str(step2),
        "reference_image_used": str(ref_path),
        "slots": trace_slots,
    })
    _write_json(manifest_path, manifest)
    return str(ref_path)


def _crop_image_to_ratio(src: str, dst: str, target_ratio: float) -> str:
    from PIL import Image

    im = Image.open(src).convert("RGB")
    ratio = im.width / im.height
    if ratio > target_ratio:
        new_w = int(im.height * target_ratio)
        left = max(0, (im.width - new_w) // 2)
        im = im.crop((left, 0, left + new_w, im.height))
    elif ratio < target_ratio:
        new_h = int(im.width / target_ratio)
        top = max(0, (im.height - new_h) // 2)
        im = im.crop((0, top, im.width, top + new_h))
    Path(dst).parent.mkdir(parents=True, exist_ok=True)
    im.save(dst, quality=95)
    return dst


def _truthy_or_auto(value: Any) -> Optional[bool]:
    if value is None:
        return None
    text = str(value).strip().lower()
    if text in {"auto", "smart", "default"}:
        return None
    return _truthy(value)


def _trim_white_margins(src: str, dst: str, threshold: int = 248, padding: int = 10) -> str:
    from PIL import Image, ImageChops

    im = Image.open(src).convert("RGB")
    bg = Image.new("RGB", im.size, (255, 255, 255))
    diff = ImageChops.difference(im, bg).convert("L")
    # Treat near-white pixels as background so screenshot/document margins are
    # removed, while retaining axes, labels, arrows, colored boxes, and text.
    mask = diff.point(lambda p: 255 if p > (255 - threshold) else 0)
    bbox = mask.getbbox()
    if not bbox:
        return src
    left, top, right, bottom = bbox
    left = max(0, left - padding)
    top = max(0, top - padding)
    right = min(im.width, right + padding)
    bottom = min(im.height, bottom + padding)
    if (right - left) >= im.width * 0.96 and (bottom - top) >= im.height * 0.96:
        return src
    Path(dst).parent.mkdir(parents=True, exist_ok=True)
    im.crop((left, top, right, bottom)).save(dst, quality=95)
    return dst


def _display_asset_for_slot(slot: Dict[str, Any], asset_path: str, output_dir: str, slide_number: int) -> str:
    explicit = _truthy_or_auto(slot.get("trim_whitespace"))
    detail_level = str(slot.get("detail_level") or "").lower()
    tags = {str(x).lower() for x in (slot.get("asset_tags") or [])}
    should_trim = explicit
    if should_trim is None:
        should_trim = detail_level in {"high", "critical"} and "white-background" in tags
    if not should_trim:
        return asset_path
    dst = os.path.join(
        output_dir,
        "assets",
        "_trimmed",
        f"slide-{slide_number:02d}-{slot.get('id', 'asset')}.jpg",
    )
    try:
        return _trim_white_margins(asset_path, dst)
    except Exception as e:
        print(f"(!)  白边裁切失败，使用原图: {asset_path} ({e})")
        return asset_path


def _add_external_image_overlays(prs: Any, slide: Any, slide_number: int, slide_spec: Dict[str, Any], output_dir: str) -> int:
    """Draw exact slot panels and add real image overlays to a PPT slide."""
    try:
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        return 0

    slide_spec = prepare_external_image_slots(slide_spec, output_dir) or slide_spec
    slots = _slots_with_real_sources(_collect_external_image_slots(slide_spec), output_dir)
    added = 0
    for idx, slot in enumerate(slots, start=1):
        asset_path = _resolve_asset_path(slot.get("source"), output_dir)
        if not asset_path:
            continue
        x, y, w, h = slot["bbox"]
        pad = max(0.0, float(slot.get("padding") or 0.0))
        left = int(x * prs.slide_width)
        top = int(y * prs.slide_height)
        box_w = int(w * prs.slide_width)
        box_h = int(h * prs.slide_height)

        if bool(slot.get("mask_placeholder")):
            print(
                f"(!)  slide {slide_number} / {slot.get('id')}: "
                "mask_placeholder 已忽略；不会在真实图片下方添加填充底图。"
            )

        # Optional final frame only. Never add a fill-only mask under the real
        # image: it becomes a selectable white rectangle in PowerPoint.
        if bool(slot.get("draw_frame")) and int(slot.get("outline_width") or 0) > 0:
            mask_bleed = max(0.0, float(slot.get("mask_bleed") or 0.0))
            mask_x = max(0.0, x - mask_bleed)
            mask_y = max(0.0, y - mask_bleed)
            mask_r = min(1.0, x + w + mask_bleed)
            mask_b = min(1.0, y + h + mask_bleed)
            frame = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                int(mask_x * prs.slide_width),
                int(mask_y * prs.slide_height),
                int((mask_r - mask_x) * prs.slide_width),
                int((mask_b - mask_y) * prs.slide_height),
            )
            outline_rgb = _hex_to_rgb(slot.get("outline"), (26, 26, 26))
            frame.fill.background()
            frame.line.color.rgb = RGBColor(*outline_rgb)
            frame.line.width = int(slot.get("outline_width") or 2) * 12700

        bleed = max(0.0, float(slot.get("bleed") or 0.0))
        inner_x = max(0.0, x + pad - bleed)
        inner_y = max(0.0, y + pad - bleed)
        inner_w = max(0.001, min(1.0 - inner_x, w - 2 * pad + 2 * bleed))
        inner_h = max(0.001, min(1.0 - inner_y, h - 2 * pad + 2 * bleed))
        pic_left = int(inner_x * prs.slide_width)
        pic_top = int(inner_y * prs.slide_height)
        pic_box_w = int(inner_w * prs.slide_width)
        pic_box_h = int(inner_h * prs.slide_height)
        box_ratio = pic_box_w / pic_box_h

        pic_path = _display_asset_for_slot(slot, asset_path, output_dir, slide_number)
        if slot.get("fit") == "cover":
            pic_path = os.path.join(
                output_dir,
                "assets",
                "_cropped",
                f"slide-{slide_number:02d}-{slot['id']}-{idx}.jpg",
            )
            try:
                pic_path = _crop_image_to_ratio(asset_path, pic_path, box_ratio)
            except Exception as e:
                print(f"(!)  cover 裁剪失败，回退 contain: {asset_path} ({e})")
                pic_path = asset_path

        try:
            from PIL import Image
            with Image.open(pic_path) as im:
                ratio = im.width / im.height
        except Exception:
            ratio = box_ratio

        if ratio >= box_ratio:
            draw_w = pic_box_w
            draw_h = int(pic_box_w / ratio)
            draw_left = pic_left
            draw_top = pic_top + int((pic_box_h - draw_h) / 2)
        else:
            draw_h = pic_box_h
            draw_w = int(pic_box_h * ratio)
            draw_left = pic_left + int((pic_box_w - draw_w) / 2)
            draw_top = pic_top
        slide.shapes.add_picture(pic_path, draw_left, draw_top, width=draw_w, height=draw_h)
        added += 1
    return added


def create_final_overlay_trace(
    output_dir: str,
    slide_number: int,
    slide_spec: Optional[Dict[str, Any]],
) -> Optional[str]:
    """Create step4 preview: generated background plus real images at exact slots."""
    slide_spec = prepare_external_image_slots(slide_spec, output_dir)
    slots = _slots_with_real_sources(_collect_external_image_slots(slide_spec), output_dir)
    if not slots:
        return None
    try:
        from PIL import Image
    except ImportError:
        return None

    bg_path = Path(output_dir) / "images" / f"slide-{slide_number:02d}.png"
    if not bg_path.is_file():
        return None

    width, height = SLIDE_TRACE_WIDTH, SLIDE_TRACE_HEIGHT
    canvas = Image.open(bg_path).convert("RGB").resize((width, height))
    trace_slots = []
    for slot in slots:
        asset_path = _resolve_asset_path(slot.get("source"), output_dir)
        if not asset_path:
            continue
        display_asset_path = _display_asset_for_slot(slot, asset_path, output_dir, slide_number)
        final_rect_px = _compute_final_image_rect_px(slot, display_asset_path, width, height)
        _paste_asset_into_rect(canvas, display_asset_path, final_rect_px, fit=slot.get("fit", "contain"))
        trace_slots.append({
            "id": slot.get("id"),
            "source": asset_path,
            "display_source": display_asset_path,
            "final_image_rect_px": final_rect_px,
            "slot_bbox_norm": slot.get("bbox"),
            "inner_rect_norm": _slot_inner_rect(slot),
            "layout_intent": slot.get("layout_intent"),
            "auto_layout_reason": slot.get("auto_layout_reason"),
        })

    trace = _trace_dir(output_dir, slide_number)
    trace.mkdir(parents=True, exist_ok=True)
    step3 = trace / "step3-generated-background.png"
    step4 = trace / "step4-final-overlay-preview.png"
    try:
        Image.open(bg_path).convert("RGB").resize((width, height)).save(step3)
    except Exception:
        pass
    canvas.save(step4)

    manifest_path = trace / "manifest.json"
    manifest = _load_trace_manifest(manifest_path)
    manifest.update({
        "slide_number": slide_number,
        "step3_generated_background": str(step3),
        "step4_final_overlay_preview": str(step4),
        "final_overlay_slots": trace_slots,
    })
    _write_json(manifest_path, manifest)
    return str(step4)


def _coerce_suggested_position(value: Any) -> Optional[List[float]]:
    if not isinstance(value, (list, tuple)) or len(value) != 4:
        return None
    try:
        nums = [float(v) for v in value]
    except (TypeError, ValueError):
        return None
    # Some vision models return pixel-ish coordinates. Normalize against the
    # trace canvas when values are clearly not 0..1.
    if any(abs(v) > 1.5 for v in nums):
        x1, y1, x2, y2 = nums
        if x2 > x1 and y2 > y1:
            nums = [x1 / SLIDE_TRACE_WIDTH, y1 / SLIDE_TRACE_HEIGHT, (x2 - x1) / SLIDE_TRACE_WIDTH, (y2 - y1) / SLIDE_TRACE_HEIGHT]
        else:
            nums = [nums[0] / SLIDE_TRACE_WIDTH, nums[1] / SLIDE_TRACE_HEIGHT, nums[2] / SLIDE_TRACE_WIDTH, nums[3] / SLIDE_TRACE_HEIGHT]
    return _parse_bbox(nums)


def _expand_position_for_readability(spec: Dict[str, Any], slot_id: str, position: List[float]) -> List[float]:
    elem = ((spec.get("elements") or {}).get(slot_id) or {})
    detail_level = _external_asset_detail_level(spec, slot_id, elem)
    if detail_level not in {"high", "critical"}:
        return position
    ratio = None
    try:
        ratio = float(elem.get("asset_ratio"))
    except (TypeError, ValueError):
        pass

    x, y, w, h = position
    min_w, min_h = (0.50, 0.68) if detail_level == "critical" else (0.46, 0.68)
    if ratio and ratio >= 2.0:
        min_w, min_h = (0.78, 0.38) if detail_level == "critical" else (0.70, 0.30)
    elif ratio and ratio < 0.85:
        min_w, min_h = (0.38, 0.78) if detail_level == "critical" else (0.30, 0.70)

    new_w = max(w, min_w)
    new_h = max(h, min_h)
    cx = x + w / 2
    cy = y + h / 2
    new_x = max(0.06, min(0.94 - new_w, cx - new_w / 2))
    new_y = max(0.08, min(0.90 - new_h, cy - new_h / 2))
    return _parse_bbox([new_x, new_y, new_w, new_h]) or position


def _overlay_review_passed(review: Optional[Dict[str, Any]]) -> bool:
    if not review:
        return False
    return review.get("ok") is True and review.get("white_placeholder_issue") is not True


def _bbox_delta(a: Optional[List[float]], b: Optional[List[float]]) -> float:
    if not a or not b:
        return 1.0
    return sum(abs(float(x) - float(y)) for x, y in zip(a, b))


def _fallback_repair_position(spec: Dict[str, Any], slot_id: str) -> Optional[List[float]]:
    """Pick a conservative slot if vision reports a problem but gives no bbox.

    Prefer right-side regions because generated text usually flows from the
    left/top in the built-in styles. Skip candidates that are effectively the
    current position so repeated repair rounds can still make progress.
    """
    elem = ((spec.get("elements") or {}).get(slot_id) or {})
    current = _parse_bbox(elem.get("position") or elem.get("computed_bbox") or elem.get("bbox"))
    asset_ratio = None
    try:
        asset_ratio = float(elem.get("asset_ratio"))
    except (TypeError, ValueError):
        pass
    detail_level = _external_asset_detail_level(spec, slot_id, elem)

    if detail_level in {"high", "critical"}:
        if detail_level == "critical" and asset_ratio and asset_ratio >= 2.0:
            candidates = [
                [0.06, 0.48, 0.88, 0.44],
                [0.08, 0.44, 0.84, 0.48],
                [0.30, 0.08, 0.64, 0.84],
            ]
        elif detail_level == "critical" and asset_ratio and asset_ratio < 0.85:
            candidates = [
                [0.52, 0.08, 0.42, 0.84],
                [0.46, 0.08, 0.48, 0.84],
                [0.10, 0.08, 0.42, 0.84],
            ]
        elif detail_level == "critical":
            candidates = [
                [0.34, 0.06, 0.62, 0.88],
                [0.30, 0.08, 0.66, 0.84],
                [0.32, 0.12, 0.63, 0.78],
                [0.35, 0.15, 0.60, 0.70],
            ]
        elif asset_ratio and asset_ratio >= 2.0:
            candidates = [
                [0.12, 0.58, 0.78, 0.32],
                [0.10, 0.54, 0.82, 0.36],
                [0.50, 0.16, 0.42, 0.56],
            ]
        elif asset_ratio and asset_ratio < 0.85:
            candidates = [
                [0.62, 0.12, 0.30, 0.72],
                [0.58, 0.14, 0.34, 0.68],
                [0.10, 0.18, 0.30, 0.68],
            ]
        else:
            candidates = [
                [0.50, 0.18, 0.42, 0.62],
                [0.48, 0.16, 0.44, 0.66],
                [0.12, 0.58, 0.78, 0.32],
            ]
    elif asset_ratio and asset_ratio < 0.85:
        candidates = [
            [0.66, 0.14, 0.26, 0.68],
            [0.58, 0.18, 0.30, 0.58],
            [0.08, 0.18, 0.26, 0.68],
        ]
    else:
        candidates = [
            [0.55, 0.30, 0.35, 0.50],
            [0.56, 0.18, 0.36, 0.52],
            [0.58, 0.12, 0.34, 0.36],
            [0.14, 0.66, 0.72, 0.24],
        ]
    for candidate in candidates:
        parsed = _parse_bbox(candidate)
        if parsed and _bbox_delta(current, parsed) >= 0.08:
            return parsed
    return None


def _region_zone(region: Optional[List[float]]) -> str:
    if not region:
        return "unknown"
    x, y, w, h = region
    cx = x + w / 2
    cy = y + h / 2
    if y >= 0.44:
        return "bottom"
    if y + h <= 0.54:
        return "top"
    if x + w <= 0.52:
        return "left"
    if x >= 0.48:
        return "right"
    if 0.34 <= cx <= 0.66 and 0.30 <= cy <= 0.70:
        return "center"
    return "mixed"


def _template_preserving_repair_position(
    spec: Dict[str, Any],
    slot_id: str,
    suggested: List[float],
) -> List[float]:
    """Keep repair close to the template's chosen media zone when possible."""
    elem = ((spec.get("elements") or {}).get(slot_id) or {})
    reason = str(elem.get("auto_layout_reason") or "")
    if "template-aware" not in reason and not spec.get("template_layout_profile"):
        return suggested
    current = _parse_bbox(elem.get("position") or elem.get("computed_bbox") or elem.get("bbox"))
    if not current:
        return suggested
    current_zone = _region_zone(current)
    suggested_zone = _region_zone(suggested)
    if current_zone == suggested_zone or current_zone in {"unknown", "mixed"}:
        return suggested

    detail_level = _external_asset_detail_level(spec, slot_id, elem)
    large = detail_level in {"high", "critical"}
    zone_candidates = {
        "left": [0.06, 0.08, 0.44, 0.84] if large else [0.06, 0.14, 0.40, 0.72],
        "right": [0.50, 0.08, 0.44, 0.84] if large else [0.54, 0.14, 0.40, 0.72],
        "bottom": [0.08, 0.44, 0.84, 0.50] if large else [0.10, 0.54, 0.80, 0.36],
        "top": [0.08, 0.08, 0.84, 0.50] if large else [0.10, 0.10, 0.80, 0.36],
        "center": [0.20, 0.12, 0.60, 0.76],
    }
    preserved = _parse_bbox(zone_candidates.get(current_zone))
    return preserved or suggested


def review_external_overlay_with_vision(output_dir: str, slide_number: int) -> Optional[Dict[str, Any]]:
    """Use VISION_* model to inspect whether real overlays cover important content."""
    trace = _trace_dir(output_dir, slide_number)
    step3 = trace / "step3-generated-background.png"
    step4 = trace / "step4-final-overlay-preview.png"
    manifest_path = trace / "manifest.json"
    if not (step3.is_file() and step4.is_file() and manifest_path.is_file()):
        return None
    try:
        sys.path.insert(0, str(SCRIPT_DIR))
        from template_analyzer import VisionClient
    except Exception as e:
        print(f"(!)  无法加载 vision 客户端，跳过贴图质检: {e}")
        return None

    try:
        client = VisionClient()
    except Exception as e:
        print(f"(!)  未配置 VISION_*，跳过贴图质检: {e}")
        return None

    manifest = _load_trace_manifest(manifest_path)
    slots = manifest.get("final_overlay_slots") or manifest.get("slots") or []
    user_text = (
        "请比较两张图：第一张是 gpt-image-2 生成的背景页，第二张是按代码坐标贴入真实图片后的最终预览。\n"
        "判断真实图片是否覆盖标题、正文、数字、图标、重要装饰或主体视觉；判断是否出现白色/浅色底图、占位框、遮罩矩形。\n"
        "白色/浅色底图只指生成背景或额外对象里出现的占位块；如果白色区域来自真实图片文件本身（例如白底图表、论文页、架构图截图），"
        "不要把它判为 white_placeholder_issue。\n"
        "还要判断真实图片本身是否足够大：如果真实图片是论文截图、表格、架构图、流程图、图表、代码或公式截图，"
        "必须能作为页面重点被清楚阅读；过小、只能当缩略图、细节不可辨，也应判为 ok=false，并给出更大但不遮挡内容的 suggested_position。\n"
        "对这类高细节真实图，宁可让页面正文更短、放到左侧窄栏，也要保证真实图占据足够大的主视觉区域。\n"
        "同时判断整体版式结构是否合理：如果最终页出现大片无意义空白、图文比例明显失衡、文字被挤到不自然角落、"
        "模板原有分栏/上下结构被破坏、装饰线条断裂或真实图位置让页面视觉重心明显怪异，也应判为 ok=false。"
        "这种情况下请给出既能保持模板结构、又能改善空白和比例的 suggested_position。\n"
        "注意：真实图片本身是矩形是正常现象，不要因此判为占位框。\n"
        f"当前 slots manifest: {json.dumps(slots, ensure_ascii=False)}\n"
        "只返回 JSON：{\n"
        '  "ok": boolean,\n'
        '  "cover_issue": string,\n'
        '  "white_placeholder_issue": boolean,\n'
        '  "layout_issue": string,\n'
        '  "blank_space_issue": boolean,\n'
        '  "recommendation": string,\n'
        '  "suggested_position": [x,y,w,h] 或 null,\n'
        '  "confidence": 0到1\n'
        "}\n"
        "如果建议位置，请优先返回 0..1 归一化坐标；若无法判断则返回 null。"
    )
    try:
        result = client.chat_json(
            system="你是严格的 PPT 视觉质检模型，只返回 JSON。",
            user_text=user_text,
            images=[str(step3), str(step4)],
            temperature=0.1,
        )
    except Exception as e:
        print(f"(!)  vision 贴图质检失败 slide {slide_number}: {e}")
        result = {
            "ok": False,
            "cover_issue": "",
            "white_placeholder_issue": False,
            "recommendation": f"vision overlay review failed: {str(e)[:240]}",
            "suggested_position": None,
            "confidence": 0.0,
            "vision_error": True,
            "slide_number": slide_number,
        }
        _write_json(trace / "overlay_review.json", result)
        return result
    if isinstance(result, list) and result and isinstance(result[0], dict):
        result = result[0]
    if not isinstance(result, dict):
        print(f"(!)  vision 贴图质检返回非对象 slide {slide_number}: {str(result)[:200]}")
        return None
    suggested = _coerce_suggested_position(result.get("suggested_position"))
    if suggested:
        result["suggested_position"] = suggested
    elif result.get("suggested_position") not in (None, "", []):
        result["suggested_position_raw"] = result.get("suggested_position")
        result["suggested_position"] = None
    result["slide_number"] = slide_number
    review_path = trace / "overlay_review.json"
    _write_json(review_path, result)
    print(
        f"🔎 slide {slide_number} overlay review: ok={result.get('ok')} "
        f"white_placeholder={result.get('white_placeholder_issue')} "
        f"position={result.get('suggested_position')}"
    )
    return result


def review_external_overlays(metadata: Dict[str, Any], output_dir: str, slide_numbers: List[int]) -> Dict[int, Dict[str, Any]]:
    reviews: Dict[int, Dict[str, Any]] = {}
    for n in slide_numbers:
        spec = _get_latest_slide_spec(metadata, n)
        if not _slots_with_real_sources(_collect_external_image_slots(spec), output_dir):
            continue
        result = review_external_overlay_with_vision(output_dir, n)
        if result:
            reviews[n] = result
    return reviews


def _first_external_slot_id(spec: Dict[str, Any]) -> Optional[str]:
    items = _external_element_items(spec)
    return items[0][0] if items else None


def sanitize_external_image_regions(
    output_dir: str,
    slide_number: int,
    slide_spec: Optional[Dict[str, Any]],
) -> Optional[str]:
    """Deterministically clear external-image regions in the generated slide.

    Image models may turn the spatial guide into a photo frame, card, or even a
    synthesized replacement photo. This fallback is opt-in via
    `sanitize_background: true`; the default path keeps image2 output untouched
    and relies on prompt/reference quality plus the final real-image overlay.
    """
    slide_spec = prepare_external_image_slots(slide_spec, output_dir)
    slots = [
        slot for slot in _slots_with_real_sources(_collect_external_image_slots(slide_spec), output_dir)
        if bool(slot.get("sanitize_background"))
    ]
    if not slots:
        return None
    try:
        from PIL import Image, ImageDraw
    except ImportError:
        return None

    img_path = Path(output_dir) / "images" / f"slide-{slide_number:02d}.png"
    if not img_path.is_file():
        return None

    trace = _trace_dir(output_dir, slide_number)
    trace.mkdir(parents=True, exist_ok=True)
    raw_path = trace / "step3-image2-raw.png"
    sanitized_path = trace / "step3-sanitized-background.png"

    img = Image.open(img_path).convert("RGB")
    if not raw_path.exists():
        img.save(raw_path)

    width, height = img.size
    draw = ImageDraw.Draw(img)
    cleanup_slots = []
    for slot in slots:
        cleanup_rect = slot.get("cleanup_rect")
        if cleanup_rect:
            rect_norm = _parse_bbox(cleanup_rect)
        else:
            asset_path = _resolve_asset_path(slot.get("source"), output_dir)
            if asset_path:
                final_rect_px = _compute_final_image_rect_px(slot, asset_path, width, height)
                rect_norm = [
                    final_rect_px[0] / width,
                    final_rect_px[1] / height,
                    final_rect_px[2] / width,
                    final_rect_px[3] / height,
                ]
            else:
                rect_norm = slot.get("bbox")
        if not rect_norm:
            continue
        x, y, w, h = rect_norm
        mask_bleed = max(0.0, float(slot.get("mask_bleed") or 0.0))
        x1 = max(0.0, x - mask_bleed)
        y1 = max(0.0, y - mask_bleed)
        x2 = min(1.0, x + w + mask_bleed)
        y2 = min(1.0, y + h + mask_bleed)
        rect_px = [
            int(round(x1 * width)),
            int(round(y1 * height)),
            int(round(x2 * width)),
            int(round(y2 * height)),
        ]
        fill_rgb = _hex_to_rgb(slot.get("mask_fill") or slot.get("fill"), (247, 247, 245))
        draw.rectangle(rect_px, fill=fill_rgb)
        cleanup_slots.append({
            "id": slot.get("id"),
            "cleanup_rect_norm": [x1, y1, x2 - x1, y2 - y1],
            "cleanup_rect_px": [
                rect_px[0],
                rect_px[1],
                max(0, rect_px[2] - rect_px[0]),
                max(0, rect_px[3] - rect_px[1]),
            ],
            "fill": "#%02X%02X%02X" % fill_rgb,
        })

    img.save(img_path)
    img.save(sanitized_path)

    manifest_path = trace / "manifest.json"
    manifest = _load_trace_manifest(manifest_path)
    manifest.update({
        "step3_image2_raw": str(raw_path),
        "step3_sanitized_background": str(sanitized_path),
        "sanitized_slide_image": str(img_path),
        "cleanup_slots": cleanup_slots,
    })
    _write_json(manifest_path, manifest)
    return str(sanitized_path)


# =============================================================================
# Image Generation
# =============================================================================

def generate_slide(
    prompt: str,
    slide_number: int,
    output_dir: str,
    reference_image_path: Optional[Union[str, List[str]]] = None,
    backend: str = "openai",
) -> Optional[str]:
    """Generate a single PPT slide image using gpt-image-2.

    backend:
      "openai" (default) -- direct /v1/images or /v1/chat calls, needs OPENAI_API_KEY
      "codex"            -- shell out to `codex exec`, reuses codex CLI auth
    """
    sys.path.insert(0, str(SCRIPT_DIR))
    if backend == "codex":
        from codex_backend import CodexImageBackend as _Backend
    else:
        from image_generator import GptImage2Generator as _Backend

    print(f"  Generating slide {slide_number} via {backend} backend ...")

    generator = _Backend(aspect_ratio="16:9")
    image_path = os.path.join(output_dir, "images", f"slide-{slide_number:02d}.png")

    scene_data = {
        "index": slide_number,
        "image_prompt": prompt,
    }
    generator.generate_scene_image(
        scene_data=scene_data,
        output_path=image_path,
        reference_image_path=reference_image_path,
    )
    print(f"  Slide {slide_number} saved: {image_path}")
    return image_path


# =============================================================================
# Output Generation
# =============================================================================

def save_prompts(output_dir: str, prompts_data: Dict[str, Any]) -> str:
    """Save all prompts to JSON file."""
    prompts_path = os.path.join(output_dir, "prompts.json")
    with open(prompts_path, "w", encoding="utf-8") as f:
        json.dump(prompts_data, f, ensure_ascii=False, indent=2)
    print(f"  Prompts saved: {prompts_path}")
    return prompts_path


def _load_existing_prompts(output_dir: str) -> Optional[Dict[str, Any]]:
    prompts_path = os.path.join(output_dir, "prompts.json")
    if not os.path.isfile(prompts_path):
        return None
    try:
        with open(prompts_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        if isinstance(data, dict):
            return data
    except Exception as e:
        print(f"(!)  读取旧 prompts.json 失败，将重建本次 prompts: {e}")
    return None


def _merge_prompts_data(existing: Optional[Dict[str, Any]], current: Dict[str, Any]) -> Dict[str, Any]:
    """Merge incremental prompt records without replacing real prompts by skip stubs."""
    if not existing or not isinstance(existing.get("slides"), list):
        return current

    merged_by_slide: Dict[int, Dict[str, Any]] = {}
    for slide in existing.get("slides", []):
        if not isinstance(slide, dict):
            continue
        try:
            n = int(slide.get("slide_number"))
        except (TypeError, ValueError):
            continue
        merged_by_slide[n] = slide

    for slide in current.get("slides", []):
        if not isinstance(slide, dict):
            continue
        try:
            n = int(slide.get("slide_number"))
        except (TypeError, ValueError):
            continue
        if str(slide.get("prompt", "")).startswith("(skipped") and n in merged_by_slide:
            merged_by_slide[n]["image_path"] = slide.get("image_path") or merged_by_slide[n].get("image_path")
            continue
        merged_by_slide[n] = slide

    merged = copy.deepcopy(current)
    merged["slides"] = [merged_by_slide[n] for n in sorted(merged_by_slide)]
    return merged


def generate_pptx(
    output_dir: str,
    slide_numbers: List[int],
    title: str = "Untitled",
    metadata: Optional[Dict[str, Any]] = None,
) -> Optional[str]:
    """把 images/slide-XX.png 打包成 16:9 .pptx，每页填满。

    如果 metadata/slide_spec 中声明了 external_image 元素，会额外：
    1. 可选地在同一个 bbox 上画代码确定性的槽位框/底色；
    2. 把 source 指向的真实图片作为独立 PPT picture object 贴入。

    注意：不要把 step4-final-overlay-preview.png 或手工合成图作为 PPT
    背景图打包；那会把真实图片烘焙进整页 PNG，PowerPoint 里无法单独
    选中/拖动真实图片。

    需要 python-pptx；如果没装就跳过并提示。
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
    except ImportError:
        print("(!)  跳过 .pptx 生成（缺 python-pptx，pip install python-pptx 后重试）")
        return None

    prs = Presentation()
    # 标准 16:9 PPT 尺寸：13.333 x 7.5 英寸（1280x720pt）
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # 完全空白布局

    img_dir = os.path.join(output_dir, "images")
    if metadata is None:
        meta_path = os.path.join(output_dir, METADATA_FILENAME)
        if os.path.isfile(meta_path):
            try:
                metadata = _load_metadata(output_dir)
            except SystemExit:
                metadata = None

    added = 0
    overlay_added = 0
    for i in slide_numbers:
        img_path = os.path.join(img_dir, f"slide-{i:02d}.png")
        if not os.path.exists(img_path):
            print(f"  跳过 slide-{i:02d}.png（文件不存在）")
            continue
        slide = prs.slides.add_slide(blank)
        # 图片填满整页（如果原图比例不是 16:9，python-pptx 默认按指定 width/height 拉伸）
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        if metadata:
            spec = _get_latest_slide_spec(metadata, i)
            if spec:
                overlay_added += _add_external_image_overlays(prs, slide, i, spec, output_dir)
                create_final_overlay_trace(output_dir, i, spec)
        added += 1

    if added == 0:
        print("(!)  没有可用图片，未生成 .pptx")
        return None

    # 文件名用 plan title（去除非法字符）
    safe_title = re.sub(r"[^\w\u4e00-\u9fff\-]+", "_", title)[:60] or "deck"
    pptx_path = os.path.join(output_dir, f"{safe_title}.pptx")
    prs.save(pptx_path)
    overlay_msg = f", {overlay_added} external image overlays" if overlay_added else ""
    print(f"  📑 PPTX generated: {pptx_path}  ({added} slides{overlay_msg})")
    return pptx_path


def resolve_editable_scene_dir(args: argparse.Namespace, output_dir: Union[str, Path]) -> Optional[Path]:
    """Resolve editable scenes only when the explicit opt-in mode is active."""
    if not getattr(args, "editable", False):
        return None
    explicit = getattr(args, "editable_scenes", None)
    candidate = Path(explicit) if explicit else Path(output_dir) / "editable_scenes"
    if not candidate.is_dir():
        raise ValueError(
            "可编辑模式缺少 scene 目录。请传 --editable-scenes DIR，"
            f"或创建 {Path(output_dir) / 'editable_scenes'}"
        )
    return candidate


# =============================================================================
# Commands (edit, rollback, ingest, list-sessions)
# =============================================================================

def cmd_list_sessions(args: argparse.Namespace) -> None:
    """List all generation sessions."""
    base_dir = args.output or str(CWD / OUTPUT_BASE_DIR)
    sessions = _find_sessions(base_dir)
    if not sessions:
        print("No sessions found.")
        return
    print(f"{'Session':<20} {'Slides':<8} Title")
    print("-" * 60)
    for s in sessions:
        print(f"{s['timestamp']:<20} {s['slide_count']:<8} {s['title']}")


def cmd_ingest_pptx(args: argparse.Namespace) -> None:
    """Ingest an external PPTX file and prepare it for editing.

    Renders the PPTX to PNGs, creates a session directory, and writes
    an initial metadata.json with placeholder slide_specs.  The Agent
    should then Read each page PNG and fill in the actual slide_specs.
    """
    pptx_path = args.ingest_pptx
    if not os.path.isfile(pptx_path):
        print(f"[X] File not found: {pptx_path}")
        sys.exit(1)

    # Render PPTX to PNGs
    sys.path.insert(0, str(SCRIPT_DIR))
    from render_template import render_pptx_to_pngs

    print(f"Ingesting: {pptx_path}")
    images_dir = render_pptx_to_pngs(pptx_path)
    print(f"Rendered to: {images_dir}")

    # Create session directory
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    if args.session:
        session_id = args.session
    else:
        stem = Path(pptx_path).stem
        session_id = f"{timestamp}_{stem}"
    session_dir = args.output or str(CWD / OUTPUT_BASE_DIR / session_id)
    session_images = os.path.join(session_dir, "images")
    os.makedirs(session_images, exist_ok=True)

    # Copy rendered PNGs as version 1 images
    import glob as _glob
    import shutil as _shutil

    png_files = sorted(_glob.glob(os.path.join(images_dir, "page-*.png")))
    slide_order: List[int] = []
    for i, png in enumerate(png_files, start=1):
        dest = os.path.join(session_images, f"slide-{i:02d}_v0001.png")
        _shutil.copy2(png, dest)
        # Also set as current
        current = os.path.join(session_images, f"slide-{i:02d}.png")
        _shutil.copy2(png, current)
        slide_order.append(i)

    # Build initial metadata with placeholder specs
    slides_meta: Dict[str, Any] = {}
    for i in slide_order:
        slides_meta[str(i)] = {
            "slide_number": i,
            "page_type": "content",
            "current_version": 1,
            "image_snapshot": f"images/slide-{i:02d}.png",
            "versions": [
                {
                    "version": 1,
                    "action": "ingest",
                    "spec": {
                        "layout": "(待 Agent 分析填充)",
                        "elements": {},
                    },
                    "prompt_file": "",
                    "image_snapshot": f"images/slide-{i:02d}_v0001.png",
                    "source_pptx": os.path.basename(pptx_path),
                }
            ],
        }

    metadata: Dict[str, Any] = {
        "version": 1,
        "title": Path(pptx_path).stem,
        "source_pptx": pptx_path,
        "slide_order": slide_order,
        "ingested_at": datetime.now().isoformat(),
        "slides": slides_meta,
    }
    _save_metadata(metadata, session_dir)

    print(f"Session created: {session_dir}")
    print(f"  {len(slide_order)} slides ingested")
    print(f"  metadata.json written with placeholder specs")
    print()
    print("Next: Agent should Read each page PNG and fill in slide_specs.")
    print("  For each slide, describe elements with type/content/position/style.")
    print("  Then use --edit to refine individual slides, or --plan to regenerate.")
    print(f"  Session timestamp: {session_id}")


def cmd_edit_slide(args: argparse.Namespace) -> None:
    """Edit a specific slide in an existing session.

    Reads the current slide_spec from metadata.json, applies element updates,
    constructs an edit prompt, and regenerates the slide with the original
    image as reference.
    """
    import shutil as _shutil

    session_dir = _resolve_session(args.session)
    metadata = _load_metadata(session_dir)
    images_dir = os.path.join(session_dir, "images")
    os.makedirs(images_dir, exist_ok=True)

    slide_key = str(args.edit)
    if slide_key not in metadata.get("slides", {}):
        print(f"[X] Slide {args.edit} not found in session {args.session}")
        sys.exit(1)

    slide_data = metadata["slides"][slide_key]
    current_version = slide_data["current_version"]
    page_type = slide_data.get("page_type", "content")

    # Get the current spec
    current_spec = _get_latest_slide_spec(metadata, args.edit)
    if not current_spec:
        print(f"[X] No spec found for slide {args.edit} (version {current_version})")
        print("    This session may not have slide_spec data. Consider regenerating.")
        sys.exit(1)

    # Apply element updates if provided
    element_updates: Dict[str, Any] = {}
    if args.element_updates:
        try:
            element_updates = json.loads(args.element_updates)
        except json.JSONDecodeError as e:
            print(f"[X] Invalid element-updates JSON: {e}")
            sys.exit(1)

    updated_spec = apply_spec_updates(current_spec, element_updates)
    updated_spec = prepare_external_image_slots(updated_spec, session_dir) or updated_spec

    # Construct edit prompt
    external_slot_repair = bool(args.external_slot_repair)
    if external_slot_repair:
        edit_prompt = args.edit_prompt or construct_external_slot_repair_prompt(updated_spec, element_updates)
    elif args.edit_prompt:
        edit_prompt = args.edit_prompt
        if not element_updates:
            print("(!) Warning: --edit-prompt used without --element-updates.")
            print("    The slide_spec in metadata will NOT be updated.")
            print("    Future structured edits may reference stale content.")
    elif element_updates:
        edit_prompt = construct_edit_prompt(current_spec, element_updates)
    else:
        print("[X] No edit prompt or element updates provided")
        sys.exit(1)

    # Get reference image (current slide image)
    ref_path = os.path.join(session_dir, slide_data.get("image_snapshot", ""))
    if not os.path.isfile(ref_path):
        ref_path = os.path.join(images_dir, f"slide-{args.edit:02d}.png")
    if not os.path.isfile(ref_path):
        print(f"(!) Reference image not found at {ref_path}, generating without reference")
        ref_path = None

    reference_images: Optional[Union[str, List[str]]] = ref_path
    asset_reference_image = None
    if external_slot_repair:
        asset_reference_image = create_asset_reference_skeleton(session_dir, args.edit, updated_spec)
        if asset_reference_image and ref_path:
            reference_images = [ref_path, asset_reference_image]
        elif asset_reference_image:
            reference_images = asset_reference_image
        else:
            print("(!) 未找到可用外部图片 source，无法生成槽位 skeleton；将只用当前页作为参考。")

    # Backup current image and build new version
    new_version = current_version + 1
    current_img = os.path.join(images_dir, f"slide-{args.edit:02d}.png")
    versioned_img = os.path.join(images_dir, f"slide-{args.edit:02d}_v{current_version:04d}.png")
    if os.path.exists(current_img):
        _shutil.copy2(current_img, versioned_img)
    _stabilize_version_snapshots(slide_data, args.edit, images_dir)

    # Save edit prompt
    prompt_rel = f"images/slide-{args.edit:02d}_v{new_version:04d}.txt"
    prompt_abs = os.path.join(session_dir, prompt_rel)
    with open(prompt_abs, "w", encoding="utf-8") as f:
        f.write(edit_prompt)

    print(f"Editing slide {args.edit} (v{current_version} -> v{new_version})")
    if element_updates:
        for elem_id, changes in element_updates.items():
            for k, v in changes.items():
                old = current_spec.get("elements", {}).get(elem_id, {}).get(k, "?")
                print(f"  {elem_id}.{k}: {old} -> {v}")
    print(f"  Reference: {reference_images or '(none)'}")
    print(f"  Edit prompt: {edit_prompt[:200]}...")

    # Generate new image
    image_path = generate_slide(
        edit_prompt, args.edit, session_dir,
        reference_image_path=reference_images,
        backend=args.backend,
    )

    if not image_path:
        print("[X] Edit generation failed")
        sys.exit(1)

    # Also save versioned copy of the new image
    versioned_new = os.path.join(images_dir, f"slide-{args.edit:02d}_v{new_version:04d}.png")
    _shutil.copy2(current_img, versioned_new)

    # Update metadata
    _add_slide_version(
        slide_data=slide_data,
        new_version=new_version,
        spec=updated_spec,
        action="edit",
        image_snapshot=f"images/slide-{args.edit:02d}_v{new_version:04d}.png",
        prompt_file=prompt_rel,
        edit_instruction=args.edit_instruction or args.edit_prompt or "",
        reference_version=current_version,
    )

    _save_metadata(metadata, session_dir)

    # Regenerate PPTX
    slide_nums = _collect_slide_numbers(metadata)
    if not args.no_pptx:
        generate_pptx(session_dir, slide_nums, title=metadata.get("title", "Untitled"), metadata=metadata)

    print(f"Slide {args.edit} updated to v{new_version}")
    print(f"  Image: {image_path}")
    print(f"  Metadata saved")


def cmd_rollback_slide(args: argparse.Namespace) -> None:
    """Rollback a slide to a previous version.

    Reads the target version's spec, regenerates from it (optionally with
    the target version's image as reference), and saves as a new version.
    """
    import shutil as _shutil

    session_dir = _resolve_session(args.session)
    metadata = _load_metadata(session_dir)
    images_dir = os.path.join(session_dir, "images")

    slide_key = str(args.rollback)
    if slide_key not in metadata.get("slides", {}):
        print(f"[X] Slide {args.rollback} not found in session {args.session}")
        sys.exit(1)

    slide_data = metadata["slides"][slide_key]
    current_version = slide_data["current_version"]
    target_version = args.to_version

    target = _get_version_info(metadata, args.rollback, target_version)
    if not target:
        print(f"[X] Version {target_version} not found for slide {args.rollback}")
        available = sorted(v.get("version", 0) for v in slide_data.get("versions", []))
        print(f"    Available versions: {available}")
        sys.exit(1)

    target_spec = target.get("spec", {})
    target_image_rel = target.get("image_snapshot", "")

    # Construct generation prompt from target spec
    # Use style from metadata (stored during generation), fallback to --style arg
    style_template = ""
    stored_style = metadata.get("style", "")
    if stored_style:
        # Try the stored path; it may be relative to CWD or SCRIPT_DIR
        for candidate_path in (stored_style, str(SCRIPT_DIR / stored_style)):
            if os.path.isfile(candidate_path):
                style_template = load_style_template(candidate_path)
                break
    if not style_template and hasattr(args, "style") and args.style:
        style_template = load_style_template(args.style)

    if target_spec.get("elements"):
        prompt = generate_prompt_from_spec(
            style_template or "按以下元素描述生成幻灯片页",
            target_spec,
            slide_data.get("page_type", "content"),
            args.rollback,
            len(metadata.get("slides", {})),
        )
    else:
        # Fallback: use stored prompt if available
        prompt_path = os.path.join(session_dir, target.get("prompt_file", ""))
        if prompt_path and os.path.isfile(prompt_path):
            with open(prompt_path, "r", encoding="utf-8") as f:
                prompt = f.read()
        elif target_image_rel:
            # Ingested slide with placeholder spec: regenerate from reference image
            ref_abs = os.path.join(session_dir, target_image_rel)
            if os.path.isfile(ref_abs):
                print(f"  Spec has no elements; using reference image as visual guide")
                prompt = (
                    "请生成一张与参考图视觉风格完全一致的幻灯片页面。\n"
                    "保持相同的布局、配色、字体和装饰元素。\n"
                    "如果参考图上有文字，保留其内容和位置。"
                )
            else:
                print(f"[X] No elements in spec and no reference image found for version {target_version}")
                sys.exit(1)
        else:
            print(f"[X] No elements in spec and no prompt file or reference image for version {target_version}")
            print(f"    This often happens with ingested slides before the Agent fills in slide_specs.")
            print(f"    Have the Agent Read each page PNG and fill in the spec elements first.")
            sys.exit(1)

    # Get reference image from target version
    ref_path = os.path.join(session_dir, target_image_rel) if target_image_rel else None
    if ref_path and not os.path.isfile(ref_path):
        print(f"(!) Reference image not found, regenerating without reference")
        ref_path = None

    # Backup current image
    new_version = current_version + 1
    current_img = os.path.join(images_dir, f"slide-{args.rollback:02d}.png")
    versioned_img = os.path.join(images_dir, f"slide-{args.rollback:02d}_v{current_version:04d}.png")
    if os.path.exists(current_img):
        _shutil.copy2(current_img, versioned_img)
    _stabilize_version_snapshots(slide_data, args.rollback, images_dir)

    # Save prompt
    prompt_rel = f"images/slide-{args.rollback:02d}_v{new_version:04d}.txt"
    prompt_abs = os.path.join(session_dir, prompt_rel)
    with open(prompt_abs, "w", encoding="utf-8") as f:
        f.write(prompt)

    print(f"Rolling back slide {args.rollback}: v{current_version} -> v{new_version} (from v{target_version})")

    image_path = generate_slide(
        prompt, args.rollback, session_dir,
        reference_image_path=ref_path,
        backend=args.backend,
    )

    if not image_path:
        print("[X] Rollback generation failed")
        sys.exit(1)

    # Save versioned copy of new image
    versioned_new = os.path.join(images_dir, f"slide-{args.rollback:02d}_v{new_version:04d}.png")
    _shutil.copy2(current_img, versioned_new)

    # Update metadata
    _add_slide_version(
        slide_data=slide_data,
        new_version=new_version,
        spec=target_spec,
        action="rollback",
        image_snapshot=f"images/slide-{args.rollback:02d}_v{new_version:04d}.png",
        prompt_file=prompt_rel,
        edit_instruction=f"Rollback to version {target_version}",
        reference_version=target_version,
    )

    _save_metadata(metadata, session_dir)

    # Regenerate PPTX
    slide_nums = _collect_slide_numbers(metadata)
    if not args.no_pptx:
        generate_pptx(session_dir, slide_nums, title=metadata.get("title", "Untitled"), metadata=metadata)

    print(f"Slide {args.rollback} rolled back (v{new_version}, based on v{target_version})")


# =============================================================================
# Main Entry Point
# =============================================================================

def create_argument_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="PPT Generator - Generate PPT images using OpenAI gpt-image-2",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Example usage:
  # Generate
  python scripts/generate_ppt.py --plan slides_plan.json --style styles/gradient-glass.md
  python scripts/generate_ppt.py --plan slides_plan.json --style styles/clean-tech-blue.md --slides 1,3,5
  python scripts/generate_ppt.py --plan slides_plan.json --style styles/gradient-glass.md --editable --editable-scenes editable_scenes/

  # Edit
  python scripts/generate_ppt.py --edit 3 --session 20240523_143052 --element-updates '{"subtitle":{"content":"新内容"}}'

  # Rollback
  python scripts/generate_ppt.py --rollback 3 --to-version 1 --session 20240523_143052

  # Ingest external PPTX
  python scripts/generate_ppt.py --ingest-pptx path/to/deck.pptx

  # List sessions
  python scripts/generate_ppt.py --list-sessions

Environment variables:
  OPENAI_BASE_URL:        Images API base URL (default: https://api.openai.com)
  OPENAI_API_KEY:         API key (required)
  GPT_IMAGE_MODEL_NAME:   Model name (default: gpt-image-2)
  GPT_IMAGE_QUALITY:      low / medium / high / auto (default: high)
""",
    )

    parser.add_argument("--plan", help="Path to slides plan JSON file")
    parser.add_argument("--style", help="Path to style template file (与模板输入二选一)")
    parser.add_argument("--output", help="Output directory path (default: outputs/TIMESTAMP)")
    parser.add_argument(
        "--template-pptx",
        help="用户的 .pptx 模板路径，启用「仿模板」模式",
    )
    parser.add_argument(
        "--template-images",
        help="模板每页 PNG 所在目录（强烈建议传，没有则只读 .pptx XML，不能跑 vision）",
    )
    parser.add_argument(
        "--template-profile",
        help="预先分析好的 TemplateProfile JSON。多模态 agent / 原生 Codex 可自己看图生成它，从而不需要 VISION_*。",
    )
    parser.add_argument(
        "--template-strict",
        action="store_true",
        help="高保真模式：把模板对应页作为 image reference 传给 gpt-image-2 出新图",
    )
    parser.add_argument(
        "--rebuild-template-cache",
        action="store_true",
        help="无视模板缓存重新跑 vision",
    )
    parser.add_argument(
        "--slides",
        help="Only generate specific slides, e.g. '1,3,5'",
    )
    parser.add_argument(
        "--concurrency",
        type=int,
        default=int(os.getenv("GPT_IMAGE_CONCURRENCY", "10")),
        help="并发请求数（默认 10，可用 GPT_IMAGE_CONCURRENCY 环境变量覆盖）",
    )
    parser.add_argument(
        "--no-pptx",
        action="store_true",
        help="不生成 .pptx 文件（默认会自动打包成 16:9 PPTX）",
    )
    parser.add_argument(
        "--editable",
        action="store_true",
        help="生成可编辑对象版 PPTX；需要 PowerPoint/Keynote/LibreOffice 回渲染，默认关闭",
    )
    parser.add_argument(
        "--editable-scenes",
        help="slide-XX.scene.json 所在目录；仅与 --editable 一起使用",
    )
    parser.add_argument(
        "--backend",
        choices=["openai", "codex"],
        default=os.getenv("GPT_IMAGE_BACKEND", "openai"),
        help="图片生成后端：openai=直调 OpenAI API（需 OPENAI_API_KEY，默认）；"
             "codex=启动本地 codex exec 子进程（非当前 Codex 原生 tool；更慢，仅作备用）",
    )

    # Edit / rollback / ingest commands
    parser.add_argument(
        "--edit",
        type=int,
        metavar="SLIDE_NUMBER",
        help="Edit a specific slide (requires --session). Use with --edit-prompt, --element-updates, or --external-slot-repair.",
    )
    parser.add_argument(
        "--session",
        help="Session timestamp or path (for --edit / --rollback). e.g. '20240523_143052'",
    )
    parser.add_argument(
        "--edit-prompt",
        help="Full edit prompt to send to gpt-image-2 for the edited slide",
    )
    parser.add_argument(
        "--external-slot-repair",
        action="store_true",
        help="Repair an edited slide by passing the current slide plus external-image slot skeleton as references.",
    )
    parser.add_argument(
        "--auto-review-overlays",
        action="store_true",
        help="After PPTX packaging, use VISION_* to review real-image overlays for collisions/placeholders.",
    )
    parser.add_argument(
        "--auto-repair-overlays",
        action="store_true",
        help="After review, automatically run bounded external-slot repair rounds for slides with overlay issues.",
    )
    parser.add_argument(
        "--edit-instruction",
        help="Human-readable description of the edit (stored in metadata)",
    )
    parser.add_argument(
        "--element-updates",
        help="JSON mapping element_id -> {key: new_value} to update the slide_spec."
             " e.g. '{\"subtitle\": {\"content\": \"新副标题\"}}'",
    )
    parser.add_argument(
        "--rollback",
        type=int,
        metavar="SLIDE_NUMBER",
        help="Rollback a slide to a previous version (requires --session --to-version)",
    )
    parser.add_argument(
        "--to-version",
        type=int,
        help="Target version number for rollback",
    )
    parser.add_argument(
        "--ingest-pptx",
        metavar="PPTX_PATH",
        help="Ingest an external PPTX file, render to PNGs, and create a session for editing",
    )
    parser.add_argument(
        "--list-sessions",
        action="store_true",
        help="List all generation sessions",
    )

    return parser


def main() -> None:
    load_skill_env()

    parser = create_argument_parser()
    args = parser.parse_args()

    if args.editable_scenes and not args.editable:
        parser.error("--editable-scenes 只能与 --editable 一起使用")
    if args.editable and args.no_pptx:
        parser.error("--editable 与 --no-pptx 不能同时使用")

    # ── 命令分发 ────────────────────────────────────────────
    if args.list_sessions:
        cmd_list_sessions(args)
        return

    if args.ingest_pptx:
        cmd_ingest_pptx(args)
        return

    if args.edit:
        if not args.session:
            parser.error("--edit requires --session")
        cmd_edit_slide(args)
        return

    if args.rollback:
        if not args.session:
            parser.error("--rollback requires --session")
        if not args.to_version:
            parser.error("--rollback requires --to-version")
        cmd_rollback_slide(args)
        return

    if args.editable:
        from editable_pptx.workflow import require_editable_render_backend
        try:
            require_editable_render_backend()
        except RuntimeError as exc:
            parser.error(str(exc))

    # ── 生成模式：必须提供 --plan ───────────────────────────
    if not args.plan:
        parser.error("必须传 --plan（生成模式），或使用 --edit / --rollback / --ingest-pptx / --list-sessions")

    # 校验：style 与 template source 至少有一个
    use_template = bool(args.template_pptx or args.template_images or args.template_profile)
    if not use_template and not args.style:
        parser.error("必须传 --style 或 --template-pptx / --template-images / --template-profile 至少其一")

    style_template = ""
    style_layout_profile: Optional[Dict[str, Any]] = None
    if args.style:
        style_path = args.style
        if not os.path.isabs(style_path):
            candidate = SCRIPT_DIR / style_path
            if candidate.exists():
                style_path = str(candidate)
        style_template = load_style_template(style_path)
        style_layout_profile = load_style_layout_profile(style_path, style_template)
    else:
        style_path = "(template-derived)"

    # 模板模式：跑 vision 拿 TemplateProfile（带缓存）
    template_profile: Optional[Dict[str, Any]] = None
    if use_template:
        sys.path.insert(0, str(SCRIPT_DIR))
        if args.template_profile:
            with open(args.template_profile, "r", encoding="utf-8") as f:
                template_profile = json.load(f)
            template_profile = _normalize_template_profile_references(template_profile, args.template_profile)
            print(f"📦 使用预分析模板 profile: {args.template_profile}")
        else:
            # 只给了 .pptx 没给 PNG -> 自动渲染到 <cwd>/template_renders/<stem>/
            if args.template_pptx and not args.template_images:
                from render_template import render_pptx_to_pngs
                print(f"🖨️  --template-images 未指定，自动渲染 {args.template_pptx}")
                args.template_images = str(render_pptx_to_pngs(args.template_pptx))
            from template_analyzer import analyze_template
            try:
                template_profile = analyze_template(
                    pptx_path=args.template_pptx,
                    images_dir=args.template_images,
                    rebuild=args.rebuild_template_cache,
                )
            except ValueError as e:
                msg = str(e)
                if "VISION_BASE_URL" in msg or "VISION_API_KEY" in msg:
                    print("[X] 模板克隆需要先获得 TemplateProfile。")
                    print("    多模态 agent / 原生 Codex：请先看 template_renders/page-*.png，生成 profile JSON 后用 --template-profile 传入。")
                    print("    纯文本 agent（如 DeepSeek）：请配置 VISION_BASE_URL / VISION_API_KEY / VISION_MODEL_NAME。")
                raise
        if not template_profile.get("layouts"):
            print("(!)  模板分析未产出 layouts（缺 --template-images？），将回退到自由风格 prompt")
            template_profile = None
            if not args.style:
                print("[X] 模板输入不可用且未提供 --style，无法构造有效生成风格。")
                print("    请提供有效 --template-profile / VISION_*，或额外传 --style 作为 fallback。")
                sys.exit(1)
        elif args.template_strict:
            missing_refs = [
                lay.get("id", f"layout-{i + 1:02d}")
                for i, lay in enumerate(template_profile.get("layouts", []))
                if not lay.get("reference_image") or not os.path.exists(str(lay.get("reference_image")))
            ]
            if missing_refs:
                print("(!)  --template-strict 已启用，但部分 layout 没有 reference_image，相关页会退化为纯 prompt 仿作。")
                print(f"    缺 reference_image 的 layout: {', '.join(missing_refs[:8])}")
    elif style_layout_profile:
        template_profile = style_layout_profile
        print(f"📚 使用内置 style layout bank: {template_profile.get('source')}")

    with open(args.plan, "r", encoding="utf-8") as f:
        slides_plan = json.load(f)

    if args.output:
        output_dir = args.output
    else:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        # 默认输出到调用者当前工作目录，而不是 skill 安装目录
        output_dir = str(CWD / OUTPUT_BASE_DIR / timestamp)

    os.makedirs(os.path.join(output_dir, "images"), exist_ok=True)
    existing_prompts_data = _load_existing_prompts(output_dir)

    slides = slides_plan["slides"]
    total_slides = len(slides)

    if args.slides:
        target_nums = set(int(x.strip()) for x in args.slides.split(","))
        slides = [s for s in slides if s.get("slide_number") in target_nums]

    selected_slide_numbers = [s["slide_number"] for s in slides]

    print("=" * 60)
    print("PPT Generator (gpt-image-2) Started")
    print("=" * 60)
    print(f"Style: {style_path}")
    if template_profile:
        print(f"Template: {template_profile.get('source')} (hash={template_profile.get('source_hash')}, "
              f"{len(template_profile.get('layouts', []))} layouts)")
        print(f"Strict mode: {args.template_strict}")
    print(f"Slides: {len(slides)} / {total_slides}")
    print(f"Output: {output_dir}")
    print(f"Concurrency: {args.concurrency}")
    print(f"Backend: {args.backend}")
    print("=" * 60)
    print()

    prompts_data: Dict[str, Any] = {
        "metadata": {
            "title": slides_plan.get("title", "Untitled Presentation"),
            "total_slides": total_slides,
            "model": os.getenv("GPT_IMAGE_MODEL_NAME", "gpt-image-2"),
            "style": style_path,
            "template": template_profile.get("source") if template_profile else None,
            "template_strict": args.template_strict if template_profile else False,
            "generated_at": datetime.now().isoformat(),
        },
        "slides": [],
    }

    # ── metadata.json（slide_spec 版本历史）──────────────────
    # If output_dir reuses an existing session, merge rather than overwrite
    existing_meta_path = os.path.join(output_dir, METADATA_FILENAME)
    if os.path.isfile(existing_meta_path):
        metadata = _load_metadata(output_dir)
        existing_order = list(metadata.get("slide_order", []))
        existing_seen = set()
        for raw in existing_order:
            try:
                existing_seen.add(int(raw))
            except (TypeError, ValueError):
                pass
        new_order = [n for n in selected_slide_numbers if int(n) not in existing_seen]
        if new_order:
            metadata["slide_order"] = existing_order + new_order
        metadata["generated_at"] = datetime.now().isoformat()
        print(f"Merging with existing metadata ({len(existing_order)} existing + {len(new_order)} new slides)")
    else:
        metadata = {
            "version": 1,
            "title": slides_plan.get("title", "Untitled Presentation"),
            "style": style_path,
            "model": os.getenv("GPT_IMAGE_MODEL_NAME", "gpt-image-2"),
            "generated_at": datetime.now().isoformat(),
            "slide_order": selected_slide_numbers,
            "slides": {},
        }

    if template_profile:
        from template_analyzer import (
            match_layout,
            assign_layouts,
            coerce_fields,
            render_prompt_from_template,
            check_layout_reuse,
        )
        # layout 复用检测：在派发任务前打出建议，让用户决定是否中断
        reuse_warnings = check_layout_reuse(slides, template_profile)
        if reuse_warnings:
            print()
            print("=" * 60)
            print("📐 Layout 复用检测（建议尽量做到 1 page : 1 layout）")
            print("=" * 60)
            for w in reuse_warnings:
                print(w)
            print("=" * 60)
            print()
        assigned_layouts = assign_layouts(slides, template_profile)
    else:
        assigned_layouts = {}

    # 收集所有待跑任务（跳过已存在的）
    pending_tasks = []
    for slide_info in slides:
        slide_number = slide_info["slide_number"]
        page_type = slide_info.get("page_type", "content")
        content_text = slide_info.get("content", "")
        raw_slide_spec = slide_info.get("slide_spec")
        matched_layout = None
        matched_layout_id = None
        if template_profile:
            matched_layout = assigned_layouts.get(slide_number) or match_layout(slide_info, template_profile)
            matched_layout_id = matched_layout.get("id") if matched_layout else None
        if matched_layout:
            raw_slide_spec = attach_template_layout_profile(raw_slide_spec, matched_layout)
        slide_spec = prepare_external_image_slots(raw_slide_spec, output_dir)
        external_slots = _slots_with_real_sources(_collect_external_image_slots(slide_spec), output_dir)
        generation_reference_images = _collect_generation_reference_images(slide_spec, output_dir)

        existing = os.path.join(output_dir, "images", f"slide-{slide_number:02d}.png")
        if os.path.exists(existing):
            print(f"Slide {slide_number}: already exists, skipping.")
            prompts_data["slides"].append({
                "slide_number": slide_number,
                "page_type": page_type,
                "content": content_text,
                "prompt": "(skipped - already exists)",
                "image_path": existing,
            })
            # Also add to metadata for skipped slides (only if not already present)
            if slide_spec and str(slide_number) not in metadata.get("slides", {}):
                metadata["slides"][str(slide_number)] = _init_slide_metadata(
                    slide_number, page_type, slide_spec,
                    prompt_file="",
                    image_path=f"images/slide-{slide_number:02d}.png",
                )
            continue

        reference_image = None
        if template_profile:
            if matched_layout is None:
                # 模板未匹配 -> 回退到 style_template
                prompt = generate_prompt(
                    style_template, page_type, content_text, slide_number, total_slides,
                    slide_spec=slide_spec,
                    output_dir=output_dir,
                )
            else:
                fields = coerce_fields(slide_info, matched_layout)
                prompt = render_prompt_from_template(
                    profile=template_profile,
                    layout=matched_layout,
                    fields=fields,
                    language_rule=LANGUAGE_FONT_RULE.strip(),
                )
                if external_slots:
                    prompt = adapt_template_prompt_for_external_slots(prompt, external_slots)
                    prompt += _format_external_slots_constraint(external_slots)
                if args.template_strict:
                    reference_image = matched_layout.get("reference_image")
        else:
            prompt = generate_prompt(
                style_template, page_type, content_text, slide_number, total_slides,
                slide_spec=slide_spec,
                output_dir=output_dir,
            )
            matched_layout_id = None

        if generation_reference_images:
            prompt += _format_generation_reference_constraint(generation_reference_images)

        asset_reference_image = create_asset_reference_skeleton(output_dir, slide_number, slide_spec)
        if generation_reference_images:
            reference_image = _merge_reference_images(reference_image, generation_reference_images)
            print(f"Slide {slide_number}: 使用 {len(generation_reference_images)} 张真实图作为生成参考（不作为 PPT 独立对象后贴）。")
        if asset_reference_image and reference_image:
            if isinstance(reference_image, list):
                reference_image = [*reference_image, asset_reference_image]
            else:
                reference_image = [reference_image, asset_reference_image]
            print(f"Slide {slide_number}: 使用模板 reference + 外部图片槽位骨架 reference。")
        elif asset_reference_image:
            reference_image = asset_reference_image

        pending_tasks.append({
            "slide_number": slide_number,
            "page_type": page_type,
            "content": content_text,
            "prompt": prompt,
            "reference_image": reference_image,
            "generation_reference_images": generation_reference_images,
            "asset_reference_image": asset_reference_image,
            "layout_id": matched_layout_id,
            "slide_spec": slide_spec,
        })

    if pending_tasks:
        from concurrent.futures import ThreadPoolExecutor, as_completed

        worker_count = max(1, min(args.concurrency, len(pending_tasks)))
        print(f"📦 派发 {len(pending_tasks)} 个任务到 {worker_count} 个并发 worker...\n")

        results: Dict[int, Optional[str]] = {}

        def _run(task):
            n = task["slide_number"]
            print(f">️  [slide {n}] start ({task['page_type']}{' / ref' if task.get('reference_image') else ''})")
            try:
                path = generate_slide(
                    task["prompt"], n, output_dir,
                    reference_image_path=task.get("reference_image"),
                    backend=args.backend,
                )
                if not path:
                    raise RuntimeError("generator returned empty image path")
                sanitize_external_image_regions(output_dir, n, task.get("slide_spec"))
                print(f"[OK] [slide {n}] done")
                return n, path
            except Exception as e:
                print(f"[X] [slide {n}] failed: {e}")
                return n, None

        with ThreadPoolExecutor(max_workers=worker_count) as executor:
            futures = [executor.submit(_run, t) for t in pending_tasks]
            for fut in as_completed(futures):
                n, path = fut.result()
                results[n] = path

        # 按原顺序写回 prompts_data，同时填充 metadata
        for task in pending_tasks:
            n = task["slide_number"]
            image_path = results.get(n)
            prompts_data["slides"].append({
                "slide_number": n,
                "page_type": task["page_type"],
                "content": task["content"],
                "layout_id": task.get("layout_id"),
                "reference_image": task.get("reference_image"),
                "generation_reference_images": task.get("generation_reference_images"),
                "asset_reference_image": task.get("asset_reference_image"),
                "prompt": task["prompt"],
                "image_path": image_path,
                "slide_spec": task.get("slide_spec"),
            })

            # metadata.json: save slide_spec + version history
            spec = task.get("slide_spec") or {}
            prompt_rel = f"images/slide-{n:02d}_v0001.txt"
            img_rel = f"images/slide-{n:02d}_v0001.png"

            if image_path:
                import shutil as _shutil
                # Save prompt text for version history
                prompt_abs = os.path.join(output_dir, prompt_rel)
                os.makedirs(os.path.dirname(prompt_abs), exist_ok=True)
                with open(prompt_abs, "w", encoding="utf-8") as pf:
                    pf.write(task["prompt"])
                versioned_initial = os.path.join(output_dir, img_rel)
                _shutil.copy2(image_path, versioned_initial)

                metadata["slides"][str(n)] = _init_slide_metadata(
                    n, task["page_type"], spec, prompt_rel, img_rel
                )
            # For failed slides, still record the spec but no image
            elif spec:
                metadata["slides"][str(n)] = _init_slide_metadata(
                    n, task["page_type"], spec, "", ""
                )

    # 按 slide_number 排序，保证 prompts.json 与播放顺序一致
    prompts_data["slides"].sort(key=lambda s: s["slide_number"])
    print()

    # Save both legacy prompts.json and new metadata.json
    prompts_data = _merge_prompts_data(existing_prompts_data, prompts_data)
    save_prompts(output_dir, prompts_data)
    _save_metadata(metadata, output_dir)

    failed_slides = sorted(
        slide["slide_number"]
        for slide in prompts_data["slides"]
        if not slide.get("image_path")
    )
    if failed_slides:
        failed_str = ", ".join(str(n) for n in failed_slides)
        print(f"[X] 生成失败，未继续产出 PPTX。失败页：{failed_str}")
        print(f"    metadata.json 已保存，修复后可 --edit 继续。")
        sys.exit(1)

    pptx_path = None
    if not args.no_pptx:
        pptx_path = generate_pptx(
            output_dir,
            _collect_slide_numbers(metadata),
            title=slides_plan.get("title", "Untitled"),
            metadata=metadata,
        )

    if not args.no_pptx and (args.auto_review_overlays or args.auto_repair_overlays):
        slide_nums = _collect_slide_numbers(metadata)
        reviews = review_external_overlays(metadata, output_dir, slide_nums)
        if args.auto_repair_overlays:
            max_repair_rounds = int(os.getenv("PPT_OVERLAY_REPAIR_ROUNDS", "3"))
            for repair_round in range(1, max_repair_rounds + 1):
                repair_updates: Dict[int, Dict[str, Any]] = {}
                for n, review in reviews.items():
                    if _overlay_review_passed(review):
                        continue
                    if review and review.get("vision_error"):
                        continue
                    spec = _get_latest_slide_spec(metadata, n)
                    slot_id = _first_external_slot_id(spec)
                    if not slot_id:
                        continue
                    vision_suggested = _coerce_suggested_position(review.get("suggested_position"))
                    suggested = vision_suggested or _fallback_repair_position(spec, slot_id)
                    if not suggested:
                        continue
                    if vision_suggested:
                        suggested = _template_preserving_repair_position(spec, slot_id, suggested)
                    if not vision_suggested:
                        suggested = _expand_position_for_readability(spec, slot_id, suggested)
                    repair_updates[n] = {
                        slot_id: {
                            "position": suggested,
                            "tailor_to_asset": False,
                            "slot_strategy": "exact",
                            "layout_intent": "manual",
                        }
                    }

                if not repair_updates:
                    break

                print()
                print("=" * 60)
                print(f"🔧 自动贴图修复 第 {repair_round}/{max_repair_rounds} 轮：{len(repair_updates)} 页")
                print("=" * 60)
                for n, updates in repair_updates.items():
                    print(f"slide {n}: {updates}")
                    cmd_args = copy.copy(args)
                    cmd_args.edit = n
                    cmd_args.session = output_dir
                    cmd_args.external_slot_repair = True
                    cmd_args.element_updates = json.dumps(updates, ensure_ascii=False)
                    cmd_args.edit_prompt = None
                    cmd_args.edit_instruction = f"auto external overlay repair round {repair_round}"
                    cmd_args.no_pptx = True
                    cmd_edit_slide(cmd_args)
                    metadata = _load_metadata(output_dir)

                pptx_path = generate_pptx(
                    output_dir,
                    _collect_slide_numbers(metadata),
                    title=slides_plan.get("title", "Untitled"),
                    metadata=metadata,
                )
                reviews = review_external_overlays(metadata, output_dir, _collect_slide_numbers(metadata))

        failed_reviews = {
            n: review for n, review in reviews.items()
            if not _overlay_review_passed(review)
        }
        if failed_reviews:
            print()
            print("[X] 贴图质检未通过，未交付为成功结果。")
            for n, review in sorted(failed_reviews.items()):
                issue = review.get("cover_issue") or review.get("recommendation") or "unknown overlay issue"
                print(f"    slide {n}: {issue}")
            print(f"    Trace: {os.path.join(output_dir, 'external_image_trace')}")
            sys.exit(1)

    editable_pptx_path = None
    editable_render_dir = None
    if args.editable and not args.no_pptx:
        try:
            scene_dir = resolve_editable_scene_dir(args, output_dir)
            from editable_pptx.workflow import build_editable_output

            editable_result = build_editable_output(
                scene_dir,
                _collect_slide_numbers(metadata),
                output_dir,
                slides_plan.get("title", "Untitled"),
            )
            editable_pptx_path = str(editable_result.pptx_path)
            editable_render_dir = str(editable_result.render_dir)
        except (ValueError, OSError, RuntimeError, json.JSONDecodeError) as exc:
            print()
            print(f"[X] 可编辑 PPTX 构建失败：{exc}")
            print(f"    普通 PPTX 与已有 scene 证据已保留在：{output_dir}")
            sys.exit(1)

    print()
    print("=" * 60)
    print("Generation Complete!")
    print("=" * 60)
    print(f"Output directory: {output_dir}")
    print(f"Metadata:    {os.path.join(output_dir, 'metadata.json')}")
    if pptx_path:
        print(f"PPTX file:   {pptx_path}")
    if editable_pptx_path:
        print(f"Editable:    {editable_pptx_path}")
    if editable_render_dir:
        print(f"Editable render: {editable_render_dir}")
    print()


if __name__ == "__main__":
    main()
