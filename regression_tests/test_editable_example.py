import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.editable_pptx.workflow import load_scene_file


def test_case05_example_is_a_real_editable_deck():
    root = ROOT / "examples" / "editable-pptx" / "case05-summer-poster"
    required = [
        "original.png",
        "clean-plate.png",
        "layers/mascot-icecream-group.png",
        "edge-check-white.png",
        "edge-check-black.png",
        "slide-01.scene.json",
        "editable.pptx",
        "rendered.png",
        "quality-report.json",
        "README.md",
    ]
    for relative in required:
        assert (root / relative).is_file(), relative

    scene = load_scene_file(root / "slide-01.scene.json")
    assert scene.slide_number == 1
    assert {element.type for element in scene.elements} >= {"native_text", "native_shape", "image_layer"}

    prs = Presentation(root / "editable.pptx")
    shapes = list(prs.slides[0].shapes)
    names = {shape.name for shape in shapes}
    assert {"clean_plate", "main_title", "product_line", "mascot_icecream_group"} <= names
    assert any(
        shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.name == "mascot_icecream_group"
        for shape in shapes
    )
    assert any(shape.has_text_frame and shape.name == "main_title" for shape in shapes)
