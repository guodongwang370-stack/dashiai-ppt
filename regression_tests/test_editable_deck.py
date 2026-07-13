import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.editable_pptx.renderer import render_editable_deck
from scripts.editable_pptx.scene import EditableScene


def _png(path: Path, color: str, mode: str = "RGB") -> Path:
    Image.new(mode, (160, 90), color).save(path)
    return path


def test_renderer_creates_ordered_multi_slide_native_deck(tmp_path):
    plate_1 = _png(tmp_path / "plate1.png", "#10152A")
    plate_2 = _png(tmp_path / "plate2.png", "#F4F7FA")
    hero = _png(tmp_path / "hero.png", "#FFAA00", "RGBA")
    scene_1 = EditableScene.from_dict(
        {
            "slide_number": 2,
            "canvas": {"width": 160, "height": 90},
            "clean_plate": str(plate_1),
            "elements": [
                {
                    "id": "title",
                    "type": "native_text",
                    "bbox_px": [20, 20, 100, 20],
                    "z_index": 20,
                    "content": "第二页",
                    "style": {"color": "#FFFFFF", "font_size_pt": 24},
                },
                {
                    "id": "banner",
                    "type": "native_shape",
                    "bbox_px": [15, 15, 110, 30],
                    "z_index": 10,
                    "style": {
                        "shape": "rounded_rectangle",
                        "fill": "#7752C8",
                        "fill_transparency": 10,
                        "line": "#FFFFFF",
                        "rotation": -2,
                    },
                },
                {
                    "id": "badge",
                    "type": "native_shape",
                    "bbox_px": [130, 10, 20, 20],
                    "z_index": 30,
                    "style": {"shape": "star_5", "fill": "#FFD324", "line": "#FFD324"},
                },
            ],
        }
    )
    scene_2 = EditableScene.from_dict(
        {
            "slide_number": 1,
            "canvas": {"width": 160, "height": 90},
            "clean_plate": str(plate_2),
            "elements": [
                {
                    "id": "hero",
                    "type": "image_layer",
                    "bbox_px": [80, 10, 60, 60],
                    "z_index": 10,
                    "asset": str(hero),
                },
                {
                    "id": "flow",
                    "type": "connector",
                    "bbox_px": [15, 65, 50, 0],
                    "z_index": 20,
                    "style": {"line": "#4B73D1", "line_width_pt": 2, "end_arrow": True},
                },
            ],
        }
    )

    output = render_editable_deck([scene_1, scene_2], tmp_path / "deck.pptx")
    prs = Presentation(output)
    assert len(prs.slides) == 2
    assert [shape.name for shape in prs.slides[0].shapes] == ["clean_plate", "hero", "flow"]
    assert [shape.name for shape in prs.slides[1].shapes] == ["clean_plate", "banner", "title", "badge"]
    assert prs.slides[0].shapes[1].shape_type == MSO_SHAPE_TYPE.PICTURE
    assert prs.slides[0].shapes[2].shape_type == MSO_SHAPE_TYPE.LINE
    assert "tailEnd" in prs.slides[0].shapes[2]._element.xml
    assert 'type="triangle"' in prs.slides[0].shapes[2]._element.xml
    assert prs.slides[1].shapes[1].shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE


def test_renderer_rejects_empty_scene_list(tmp_path):
    import pytest

    with pytest.raises(ValueError, match="至少需要一个"):
        render_editable_deck([], tmp_path / "empty.pptx")
