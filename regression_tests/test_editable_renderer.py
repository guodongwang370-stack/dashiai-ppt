import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.editable_pptx.renderer import render_editable_pptx
from scripts.editable_pptx.scene import EditableScene


def _write_image(path: Path, color) -> Path:
    Image.new("RGB", (160, 90), color).save(path)
    return path


def _sample_scene(tmp_path: Path):
    clean_plate = _write_image(tmp_path / "clean.png", (10, 15, 35))
    visual = _write_image(tmp_path / "visual.png", (200, 80, 40))
    return {
        "slide_number": 1,
        "canvas": {"width": 160, "height": 90},
        "clean_plate": str(clean_plate),
        "elements": [
            {
                "id": "hero",
                "type": "image_layer",
                "bbox_px": [96, 18, 48, 54],
                "z_index": 10,
                "asset": str(visual),
            },
            {
                "id": "title",
                "type": "native_text",
                "bbox_px": [12, 18, 80, 18],
                "z_index": 20,
                "content": "年度战略复盘",
                "style": {
                    "font_face": "PingFang SC",
                    "font_size_pt": 24,
                    "font_weight": 700,
                    "color": "#FFFFFF",
                    "align": "left",
                },
            },
        ],
    }


def test_renderer_creates_native_text_and_picture_layers(tmp_path):
    scene = EditableScene.from_dict(_sample_scene(tmp_path))

    path = render_editable_pptx(scene, tmp_path / "editable.pptx")

    shapes = list(Presentation(path).slides[0].shapes)
    assert any(shape.has_text_frame and shape.text == "年度战略复盘" for shape in shapes)
    assert sum(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in shapes) == 2
    assert {shape.name for shape in shapes} >= {"clean_plate", "hero", "title"}


def test_scene_rejects_duplicate_ids(tmp_path):
    data = _sample_scene(tmp_path)
    data["elements"].append(dict(data["elements"][0]))

    try:
        EditableScene.from_dict(data)
    except ValueError as exc:
        assert "重复" in str(exc)
    else:
        raise AssertionError("duplicate element IDs should fail")


def test_scene_rejects_out_of_canvas_bbox(tmp_path):
    data = _sample_scene(tmp_path)
    data["elements"][0]["bbox_px"] = [150, 10, 20, 20]

    try:
        EditableScene.from_dict(data)
    except ValueError as exc:
        assert "bbox" in str(exc)
    else:
        raise AssertionError("out-of-canvas bbox should fail")
